"""PPTX 렌더러 — 정본 ReportModel 을 Block 단위로 순회하며 python-pptx 슬라이드 생성.

pdf_renderer 와 '같은 순서·같은 Block kind 분기'를 미러링한다:
    표지(ReportMeta) → (핵심요약) → 섹션들(블록 순회) → 면책.

씨앗은 market_report_service.to_pptx(시장조사 하드코딩)였다. 그 안의 헬퍼
(title_slide/header_bar/brand_footer/text_slide/table_slide/chart_slide/kv_table_slide/
네이티브 add_chart)를 '시장rep 전용'에서 '어떤 보고서든 Block 으로 그리는 일반형'으로 승격했다.

순수 모듈 규칙: python-pptx + 표준 라이브러리만 쓴다(reportlab/FastAPI/DB 임포트 금지).
가짜값 금지: 값 문자열화는 언제나 model.fmt_value 를 거치고, 무자료는 '데이터 없음'/'—'로 정직 표기.
"""

from __future__ import annotations

import contextlib
import io
import math
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from . import tokens as T
from .model import ReportModel, Section, fmt_value


# ── 색 변환 헬퍼: PRDS HEX → python-pptx RGBColor ──────────────────────
def _rgb(hex_str: str) -> RGBColor:
    """'#0e7490' → RGBColor(0x0e,0x74,0x90). 실패하면 딥틸로 폴백."""
    try:
        r, g, b = T.hex_to_rgb(hex_str)
        return RGBColor(r, g, b)
    except Exception:  # noqa: BLE001  # 잘못된 색값이면 브랜드색으로 안전 폴백
        r, g, b = T.hex_to_rgb(T.BRAND)
        return RGBColor(r, g, b)


# 자주 쓰는 색을 미리 만들어 둔다(토큰 → RGBColor 1:1).
BRAND = _rgb(T.BRAND)
INK = _rgb(T.INK)
MUTED = _rgb(T.MUTED)
WHITE = _rgb(T.WHITE)
ZEBRA = _rgb(T.ZEBRA)
PANEL = _rgb(T.PANEL)
LINE = _rgb(T.LINE)
AMBER = _rgb(T.AMBER)
DANGER = _rgb(T.SIGNAL["danger"])
SLATE = _rgb("#94a3b8")   # 표지 보조 텍스트(어두운 배경 위 회청색)

FONT = T.FONT_KR_GOTHIC   # '맑은 고딕' — 뷰어 로컬 한글 고딕(EastAsian)


# ── 16:9 캔버스·여백(모두 인치 단위, T.PPTX 에서 가져옴) ────────────────
_W = T.PPTX["w_in"]                    # 슬라이드 폭 13.33"
_H = T.PPTX["h_in"]                    # 슬라이드 높이 7.5"
_SAFE = T.PPTX["safe_in"]             # 좌우 안전 여백 0.6"
_HEADER_H = T.PPTX["header_in"]      # 상단 헤더바 높이 1.1"
_FOOTER_Y = T.PPTX["footer_y_in"]   # 하단 푸터 y 7.1"
_FOOTER_H = T.PPTX["footer_h_in"]   # 푸터 높이 0.4"

_CONTENT_W = _W - 2 * _SAFE          # 본문 폭(좌우 여백 제외)
_CONTENT_TOP = _HEADER_H + 0.25     # 본문 시작 y(헤더 아래)
_CONTENT_BOTTOM = _FOOTER_Y - 0.15  # 본문 끝 y(푸터 위)

# 표/블록 배치용 기본 치수
_ROW_H = 0.34    # 표 한 행 높이(인치)
_HDR_H = 0.40    # 표 헤더행 높이(인치)
_GAP = 0.14      # 블록 사이 간격(인치)


# ── 저수준 텍스트/도형 헬퍼 ────────────────────────────────────────────
def _apply_ea(run, name: str) -> None:
    """run 에 EastAsian(한글) 폰트를 지정한다.

    python-pptx 의 run.font.name 은 라틴(a:latin)만 바꾼다. 한글은 a:ea 로 따로 지정해야
    맑은 고딕이 실제로 적용된다. rPr(run 속성) 안에 a:ea 를 만들어 라틴 뒤에 끼워 넣는다.
    """
    try:
        rpr = run._r.get_or_add_rPr()
        for el in rpr.findall(qn("a:ea")):   # 중복 방지: 기존 ea 제거
            rpr.remove(el)
        ea = rpr.makeelement(qn("a:ea"), {})
        ea.set("typeface", name)
        latin = rpr.find(qn("a:latin"))
        if latin is not None:                # 스키마상 ea 는 latin 바로 뒤
            latin.addnext(ea)
        else:
            rpr.append(ea)
    except Exception:  # noqa: BLE001  # ea 지정 실패해도 문서 생성은 계속(라틴명은 이미 지정됨)
        pass


def _fmt_run(run, size: float, bold: bool, color: RGBColor, name: str = FONT) -> None:
    """run 하나의 폰트(크기·굵기·색·한글폰트)를 한 번에 지정."""
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.name = name
    f.color.rgb = color
    _apply_ea(run, name)


def _p(para, text: Any, size: float, color: RGBColor, *, bold: bool = False,
       align=PP_ALIGN.LEFT, space_after: float = 6) -> None:
    """이미 있는 문단(paragraph)에 텍스트 run 을 채우고 스타일 지정."""
    para.alignment = align
    para.space_after = Pt(space_after)
    run = para.add_run()
    run.text = fmt_value(text)
    _fmt_run(run, size, bold, color)


def _fill(shape, rgb: RGBColor, line_rgb: RGBColor | None = None) -> None:
    """도형 배경 채움 + (선택)테두리. 그림자는 제거(PRDS: 장식 최소)."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(0.75)
    with contextlib.suppress(Exception):  # 그림자 속성이 없어도 무시
        shape.shadow.inherit = False


def _style_cell(cell, text: Any, *, size: float, bold: bool, color: RGBColor,
                align=PP_ALIGN.LEFT, fill: RGBColor | None = None) -> None:
    """표 셀 하나: 배경·가운데정렬·여백·텍스트 스타일을 지정."""
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(5)
    cell.margin_right = Pt(5)
    cell.margin_top = Pt(1)
    cell.margin_bottom = Pt(1)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()            # 빈 셀 문단에 run 을 직접 추가(가짜 텍스트 없음)
    run.text = fmt_value(text)
    _fmt_run(run, size, bold, color)


def _est_text_h(text: Any, chars_per_line: int, pt: float, gap: float = 0.06) -> float:
    """긴 텍스트가 몇 인치 높이를 차지할지 보수적으로 추정(페이지 넘김 판단용).

    한글은 대략 글자폭≈글자크기라 한 줄에 들어가는 글자 수를 과소평가(=높이 과대)해서
    푸터를 침범하지 않도록 안전하게 잡는다.
    """
    s = str(text)
    hard = s.count("\n") + 1
    wrap = max(1, math.ceil(len(s) / max(1, chars_per_line)))
    lines = max(hard, wrap)
    return lines * (pt * 1.5 / 72.0) + gap


# ── 덱(Deck) 빌더: 슬라이드·본문 커서 상태를 들고 블록을 그린다 ──────────
class _Deck:
    """python-pptx Presentation 을 감싸 '본문 커서(y)' 를 관리하는 빌더.

    한 섹션을 그리다 세로 공간이 모자라면 자동으로 '(계속)' 슬라이드를 만들어 이어 그린다.
    (pdf 는 reportlab 이 자동 페이지 넘김을 하지만, pptx 는 고정 슬라이드라 직접 관리한다.)
    """

    def __init__(self, prs, meta, disclaimer_text: str):
        self.prs = prs
        self.meta = meta
        self.disclaimer_text = disclaimer_text
        self.slide = None          # 현재 그리는 슬라이드
        self.y = _CONTENT_TOP      # 현재 세로 커서(인치)
        self.section_title = ""    # 현재 섹션 제목(계속 슬라이드에 재사용)
        # 도메인 강조색(없으면 PRDS 딥틸). 헤더바·푸터·표헤더에 일관 사용.
        self.accent = _rgb(meta.accent_color) if getattr(meta, "accent_color", None) else BRAND

    # ── 슬라이드 생성 / 헤더 / 푸터 ──
    def _new_slide(self, title: str):
        """빈(Blank) 레이아웃 슬라이드 하나 + 헤더바 + 하단 브랜드 푸터. 커서 초기화."""
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide = s
        self.y = _CONTENT_TOP
        self._header_bar(title)
        self._footer()
        return s

    def _header_bar(self, title: str) -> None:
        """상단 딥틸 헤더바 + 섹션 제목(흰 굵게)."""
        bar = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(_W), Inches(_HEADER_H))
        _fill(bar, self.accent)
        tf = bar.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(_SAFE)
        tf.margin_right = Inches(_SAFE)
        _p(tf.paragraphs[0], title, 24, WHITE, bold=True, space_after=0)

    def _footer(self) -> None:
        """하단 브랜드 푸터바(매 슬라이드): 좌=브랜딩·대외비, 우=문서번호·생성일."""
        bar = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(_FOOTER_Y), Inches(_W), Inches(_FOOTER_H))
        _fill(bar, self.accent)
        tf = bar.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(_SAFE)
        tf.margin_right = Inches(_SAFE)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
        left = T.BRANDING
        if self.meta.confidential:
            left += "   |   " + T.CONFIDENTIAL_LABEL
        _p(tf.paragraphs[0], left, 9, WHITE, space_after=0)
        # 우측 메타(문서번호·생성일) — 값이 있을 때만.
        right_bits = [x for x in [
            self.meta.doc_no,
            (str(self.meta.generated_at)[:16] if self.meta.generated_at else None),
        ] if x]
        if right_bits:
            rtb = self.slide.shapes.add_textbox(
                Inches(_W - 4.6), Inches(_FOOTER_Y), Inches(4.0), Inches(_FOOTER_H))
            rtf = rtb.text_frame
            rtf.vertical_anchor = MSO_ANCHOR.MIDDLE
            _p(rtf.paragraphs[0], "  ·  ".join(right_bits), 9, WHITE, align=PP_ALIGN.RIGHT, space_after=0)

    # ── 공간 확보(넘치면 계속 슬라이드) ──
    def _ensure(self, need: float) -> None:
        """현재 커서에서 need(인치)만큼 그릴 공간이 없으면 '(계속)' 슬라이드로 넘어간다."""
        if self.slide is None:
            self._new_slide(self.section_title or "")
        elif self.y + need > _CONTENT_BOTTOM:
            self._new_slide(self.section_title + " (계속)")

    # ── 표지(ReportMeta) ──
    def _cover(self) -> None:
        """표지 슬라이드: 문서유형·부제·주소·생성일·대외비·데이터채움도·면책 + 딥틸 밴드."""
        m = self.meta
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # 잉크색 풀블리드 배경 + 상단 딥틸 밴드(브랜드 위계)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(_W), Inches(_H))
        _fill(bg, INK)
        band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.05), Inches(_W), Inches(0.12))
        _fill(band, self.accent)
        # 상단 브랜드
        bt = s.shapes.add_textbox(Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.6)).text_frame
        _p(bt.paragraphs[0], T.BRANDING, 16, self.accent, bold=True, space_after=0)
        # 타이틀 블록(문서유형 크게 → 부제/주소/생성일/대외비/채움도)
        tf = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.6), Inches(3.4)).text_frame
        tf.word_wrap = True
        _p(tf.paragraphs[0], m.title, 40, WHITE, bold=True, space_after=10)
        if m.subtitle:
            _p(tf.add_paragraph(), m.subtitle, 18, SLATE, space_after=6)
        if m.project_address:
            _p(tf.add_paragraph(), m.project_address, 20, WHITE, space_after=6)
        meta_bits = [x for x in [
            (f"생성 {m.generated_at}" if m.generated_at else None),
            (f"문서번호 {m.doc_no}" if m.doc_no else None),
        ] if x]
        if meta_bits:
            _p(tf.add_paragraph(), "  ·  ".join(meta_bits), 14, SLATE, space_after=6)
        if m.confidential:
            _p(tf.add_paragraph(), T.CONFIDENTIAL_LABEL, 14, DANGER, bold=True, space_after=4)
        if m.completeness and m.completeness.get("pct") is not None:
            _p(tf.add_paragraph(), f"데이터 채움도: {m.completeness.get('pct')}% (정직 표기)", 12, SLATE, space_after=0)
        # 하단 면책 고지
        dtf = s.shapes.add_textbox(Inches(0.9), Inches(6.5), Inches(11.6), Inches(0.85)).text_frame
        dtf.word_wrap = True
        _p(dtf.paragraphs[0], self.disclaimer_text, 9, SLATE, space_after=0)

    # ── 섹션(제목 슬라이드 시작 + 블록 순회) ──
    def _section(self, sec: Section) -> None:
        heading = f"{sec.section_no}. {sec.title}" if sec.section_no else sec.title
        self._new_slide(heading)
        self.section_title = heading
        for block in sec.blocks:
            self._block(block)

    # ── Block kind 별 렌더 디스패치(pdf_renderer._render_block 미러) ──
    def _block(self, block: Any) -> None:
        kind = getattr(block, "kind", None)

        if kind == "kv":
            if block.title:
                self._block_title(block.title)
            self._kv_table(block.rows)

        elif kind == "table":
            if block.title:
                self._block_title(block.title)
            self._data_table(block.headers, block.rows,
                             numeric_cols=block.numeric_cols, total_row=block.total_row)
            if block.caption:
                self._caption(block.caption)

        elif kind == "kpi":
            self._kpi(block.tiles)

        elif kind == "chart":
            self._chart(block)

        elif kind == "narrative":
            self._narrative(block)

        elif kind == "evidence":
            self._evidence(block)

        elif kind == "checklist":
            if block.title:
                self._block_title(block.title)
            rows = []
            for label, status in block.items or []:
                mark = "✓" if status is True else ("—" if status in (False, None) else fmt_value(status))
                rows.append([label, mark])
            if rows:
                self._data_table(["항목", "상태"], rows)
            else:
                self._caption("데이터 없음")

        elif kind == "grade":
            self._grade(block)

        elif kind == "image":
            self._image(block)

        elif kind == "disclaimer":
            h = _est_text_h(block.text, 100, 9)
            self._ensure(h)
            self._textbox(_SAFE, self.y, _CONTENT_W, h, block.text, 9, MUTED)
            self.y += h + 0.04

    # ── 소제목 / 캡션 / 일반 텍스트박스 ──
    def _block_title(self, text: Any) -> None:
        """블록 소제목(딥틸 굵게)."""
        self._ensure(0.42)
        self._textbox(_SAFE, self.y, _CONTENT_W, 0.34, text, 14, self.accent, bold=True)
        self.y += 0.38

    def _caption(self, text: Any) -> None:
        """캡션/출처(작은 회색)."""
        h = _est_text_h(text, 100, 10)
        self._ensure(h)
        self._textbox(_SAFE, self.y, _CONTENT_W, h, text, 10, MUTED)
        self.y += h + 0.04

    def _textbox(self, left: float, top: float, width: float, height: float,
                 text: Any, size: float, color: RGBColor, *, bold: bool = False,
                 align=PP_ALIGN.LEFT):
        """한 문단짜리 텍스트박스를 현재 슬라이드에 추가."""
        tb = self.slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
        _p(tf.paragraphs[0], text, size, color, bold=bold, align=align, space_after=0)
        return tb

    # ── KV 2열표(헤더 없이 라벨열 음영) ──
    def _kv_table(self, rows) -> None:
        rows = list(rows or [])
        if not rows:                       # 무자료 정직 표기
            self._ensure(0.4)
            self._textbox(_SAFE, self.y, _CONTENT_W, 0.32, T.EMPTY_MARK, 11, MUTED)
            self.y += 0.32 + _GAP
            return
        n = len(rows)
        idx = 0
        while idx < n:
            avail = _CONTENT_BOTTOM - self.y
            cap = int(avail // _ROW_H)
            if cap < 1:                    # 이 슬라이드엔 한 행도 못 들어감 → 계속 슬라이드
                self._new_slide(self.section_title + " (계속)")
                cap = max(1, int((_CONTENT_BOTTOM - self.y) // _ROW_H))
            chunk = rows[idx: idx + cap]
            nr = len(chunk)
            gf = self.slide.shapes.add_table(
                nr, 2, Inches(_SAFE), Inches(self.y), Inches(_CONTENT_W), Inches(nr * _ROW_H))
            tbl = gf.table
            self._plain_table(tbl)
            tbl.columns[0].width = Inches(3.7)               # 라벨열
            tbl.columns[1].width = Inches(_CONTENT_W - 3.7)  # 값열
            for i, kv in enumerate(chunk):
                k = kv[0] if len(kv) > 0 else ""
                v = kv[1] if len(kv) > 1 else ""
                tbl.rows[i].height = Inches(_ROW_H)
                _style_cell(tbl.cell(i, 0), k, size=11, bold=True, color=INK, fill=ZEBRA)
                _style_cell(tbl.cell(i, 1), v, size=11, bold=False, color=INK, fill=WHITE)
            self.y += nr * _ROW_H + _GAP
            idx += nr
            if idx < n:
                self._new_slide(self.section_title + " (계속)")

    # ── 일반 데이터표(딥틸 헤더 + zebra + 숫자열 우측정렬 + 합계행) ──
    def _data_table(self, headers, rows, *, numeric_cols=(), total_row: bool = False) -> None:
        headers = list(headers or [])
        rows = list(rows or [])
        ncol = max(1, len(headers))
        if not rows:                       # 무자료 정직 표기
            self._ensure(0.4)
            self._textbox(_SAFE, self.y, _CONTENT_W, 0.32, "데이터 없음", 11, MUTED)
            self.y += 0.32 + _GAP
            return
        n = len(rows)
        idx = 0
        while idx < n:
            avail = _CONTENT_BOTTOM - self.y - _HDR_H
            cap = int(avail // _ROW_H)
            if cap < 1:                    # 헤더+한 행도 안 들어감 → 계속 슬라이드
                self._new_slide(self.section_title + " (계속)")
                cap = max(1, int((_CONTENT_BOTTOM - self.y - _HDR_H) // _ROW_H))
            chunk = rows[idx: idx + cap]
            chunk_has_total = total_row and (idx + len(chunk) >= n)
            self._draw_table_chunk(headers, chunk, ncol, list(numeric_cols or []), chunk_has_total)
            self.y += _HDR_H + len(chunk) * _ROW_H + _GAP
            idx += len(chunk)
            if idx < n:
                self._new_slide(self.section_title + " (계속)")

    def _draw_table_chunk(self, headers, chunk, ncol, numeric_cols, total_flag) -> None:
        nrows = 1 + len(chunk)
        total_h = _HDR_H + len(chunk) * _ROW_H
        gf = self.slide.shapes.add_table(
            nrows, ncol, Inches(_SAFE), Inches(self.y), Inches(_CONTENT_W), Inches(total_h))
        tbl = gf.table
        self._plain_table(tbl)
        col_w = Inches(_CONTENT_W / ncol)
        for c in range(ncol):
            tbl.columns[c].width = col_w
        tbl.rows[0].height = Inches(_HDR_H)
        for r in range(1, nrows):
            tbl.rows[r].height = Inches(_ROW_H)
        # 헤더(딥틸 배경 · 흰 굵게)
        for c in range(ncol):
            h = headers[c] if c < len(headers) else ""
            _style_cell(tbl.cell(0, c), h, size=11, bold=True, color=WHITE, fill=self.accent)
        # 본문(행 교대 zebra · 숫자열 우측정렬 · 합계행은 굵게+패널배경)
        for ri, row in enumerate(chunk):
            r = ri + 1
            is_total = total_flag and (ri == len(chunk) - 1)
            zebra = ZEBRA if (r % 2 == 0) else WHITE
            fill = PANEL if is_total else zebra
            for c in range(ncol):
                val = row[c] if c < len(row) else ""
                align = PP_ALIGN.RIGHT if c in numeric_cols else PP_ALIGN.LEFT
                _style_cell(tbl.cell(r, c), val, size=11, bold=is_total, color=INK, align=align, fill=fill)

    @staticmethod
    def _plain_table(tbl) -> None:
        """기본 표 스타일(굵은 첫행·자동 줄무늬)을 끈다 → 우리가 셀별로 직접 칠한다."""
        try:
            tbl.first_row = False
            tbl.horz_banding = False
        except Exception:  # noqa: BLE001  # 속성 미지원이면 무시(셀 채움으로 덮어씀)
            pass

    # ── KPI 타일(둥근 사각형 3~4개) ──
    def _kpi(self, tiles) -> None:
        tiles = list(tiles or [])
        if not tiles:
            return
        per_row = 4 if len(tiles) >= 4 else len(tiles)   # 한 줄 3~4개
        gap = 0.22
        th = 1.3
        for i in range(0, len(tiles), per_row):
            group = tiles[i: i + per_row]
            nn = len(group)
            self._ensure(th + 0.05)
            tw = (_CONTENT_W - gap * (nn - 1)) / nn
            for j, tile in enumerate(group):
                self._tile(_SAFE + j * (tw + gap), self.y, tw, th, tile)
            self.y += th + _GAP

    def _tile(self, left: float, top: float, w: float, h: float, tile) -> None:
        """KPI 타일 1개: 둥근 사각형 + 라벨(작게)/수치(크게·신호색)/기준(작게)."""
        shp = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
        _fill(shp, WHITE, line_rgb=LINE)
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(6)
        # 라벨(위, 회색)
        _p(tf.paragraphs[0], tile.label, 10, MUTED, align=PP_ALIGN.CENTER, space_after=2)
        # 수치(크게) — signal(hex) 있으면 그 색으로, 없으면 잉크색
        vcolor = _rgb(tile.signal) if (tile.signal and str(tile.signal).startswith("#")) else INK
        pv = tf.add_paragraph()
        _p(pv, tile.value, 24, vcolor, bold=True, align=PP_ALIGN.CENTER, space_after=2)
        # 기준(아래, 회색) — 있을 때만
        if tile.basis:
            pb = tf.add_paragraph()
            _p(pb, tile.basis, 9, MUTED, align=PP_ALIGN.CENTER, space_after=0)

    # ── 차트(bar/line/pie=네이티브, waterfall/tornado=표 폴백) ──
    def _chart(self, block) -> None:
        if block.chart_type in ("bar", "line", "pie"):
            self._native_chart(block)
        else:
            # 벡터로 못 그리는 waterfall/tornado 는 표로 정직 폴백(가짜 곡선 금지).
            if block.title:
                self._block_title(block.title)
            if block.series:
                headers = ["구분", *[s.name for s in block.series]]
                rows = []
                for i, cat in enumerate(block.categories):
                    row = [cat]
                    for s in block.series:
                        row.append(fmt_value(s.values[i]) if i < len(s.values) else T.EMPTY_MARK)
                    rows.append(row)
                self._data_table(headers, rows, numeric_cols=list(range(1, len(headers))))
            else:
                self._caption("데이터 없음")
        if block.caption:
            self._caption(block.caption)

    def _native_chart(self, block) -> None:
        """python-pptx 네이티브 차트(편집 가능한 벡터). 씨앗의 CategoryChartData 패턴을 일반화."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

        xmap = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
        }
        xtype = xmap[block.chart_type]
        ch_h = 4.2
        self._ensure(ch_h)   # 차트는 크므로 공간 없으면 다음 슬라이드로

        def _nums(vals):     # 차트는 숫자만 받는다(None→0). 값 없으면 0으로.
            return tuple((v if isinstance(v, (int, float)) else 0) for v in (vals or []))

        cd = CategoryChartData()
        cd.categories = [str(c) for c in (block.categories or [""])]
        if block.chart_type == "pie":
            ser = block.series[0] if block.series else None
            cd.add_series(ser.name if ser else "", _nums(ser.values) if ser else (0,))
        elif block.series:
            for s in block.series:
                cd.add_series(s.name, _nums(s.values))
        else:
            cd.add_series("", (0,))

        gf = self.slide.shapes.add_chart(
            xtype, Inches(_SAFE), Inches(self.y), Inches(_CONTENT_W), Inches(ch_h), cd)
        chart = gf.chart

        # 차트 전체 폰트(한글 포함)
        try:
            chart.font.size = Pt(10)
            chart.font.name = FONT
        except Exception:  # noqa: BLE001
            pass
        # 제목
        try:
            chart.has_title = True
            chart.chart_title.text_frame.text = fmt_value(block.title)
            for para in chart.chart_title.text_frame.paragraphs:
                for run in para.runs:
                    _fmt_run(run, 12, True, INK)
        except Exception:  # noqa: BLE001
            pass
        # 범례(선/파이 또는 다계열 막대만)
        try:
            multi = bool(block.series) and len(block.series) > 1
            if block.chart_type in ("line", "pie") or multi:
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
            else:
                chart.has_legend = False
        except Exception:  # noqa: BLE001
            pass
        # 세로축 제목
        try:
            if block.y_axis_label and block.chart_type in ("bar", "line"):
                va = chart.value_axis
                va.has_title = True
                va.axis_title.text_frame.text = fmt_value(block.y_axis_label)
        except Exception:  # noqa: BLE001
            pass
        # 시리즈/조각 색(PRDS 계열색)
        try:
            if block.chart_type == "pie":
                pts = chart.plots[0].series[0].points
                for i, pt in enumerate(pts):
                    col = _rgb(T.SERIES_COLORS[i % len(T.SERIES_COLORS)])
                    pt.format.fill.solid()
                    pt.format.fill.fore_color.rgb = col
            else:
                for i, ser in enumerate(chart.series):
                    col = _rgb(T.SERIES_COLORS[i % len(T.SERIES_COLORS)])
                    if block.chart_type == "line":
                        ser.format.line.color.rgb = col
                    else:
                        ser.format.fill.solid()
                        ser.format.fill.fore_color.rgb = col
        except Exception:  # noqa: BLE001  # 색 지정 실패해도 차트 자체는 정상
            pass

        self.y += ch_h + _GAP

    # ── 서술(문단) ──
    def _narrative(self, block) -> None:
        if block.title:
            self._block_title(block.title)
        for para in block.paragraphs or []:
            text = str(para).strip()
            if not text:
                continue
            h = _est_text_h(text, 48, 15)
            self._ensure(h)
            self._textbox(_SAFE, self.y, _CONTENT_W, h, text, 15, INK)
            self.y += h + 0.06

    # ── 근거(각 항목 한 줄, 저신뢰는 앰버) ──
    def _evidence(self, block) -> None:
        if block.title:
            self._block_title(block.title)
        for ev in block.items or []:
            conf = str(ev.confidence).lower() if ev.confidence else ""
            low = conf in ("low", "med", "medium")
            color = AMBER if low else MUTED          # R4: 저신뢰(low/med)는 앰버로 명확히
            # 한 줄 구성: 값(굵게) · 근거 · 출처 · (신뢰도) · 법령
            segs: list[tuple[str, bool]] = [(fmt_value(ev.value), True)]
            if ev.basis:
                segs.append((f"근거: {fmt_value(ev.basis)}", False))
            if ev.source:
                segs.append((f"출처: {fmt_value(ev.source)}", False))
            if low:
                segs.append((f"(신뢰도 {fmt_value(ev.confidence)})", False))
            # 법령 링크는 verified(high 또는 미표기)만 노출(pdf 규칙 동일)
            if ev.legal_link and (conf == "high" or not ev.confidence):
                segs.append((f"법령: {fmt_value(ev.legal_link)}", False))
            joined = " · ".join(s for s, _ in segs)
            h = _est_text_h(joined, 84, 10)
            self._ensure(h)
            tb = self.slide.shapes.add_textbox(
                Inches(_SAFE), Inches(self.y), Inches(_CONTENT_W), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(2)
            tf.margin_right = Pt(2)
            tf.margin_top = Pt(1)
            tf.margin_bottom = Pt(1)
            p = tf.paragraphs[0]
            for k, (seg, bold) in enumerate(segs):
                run = p.add_run()
                run.text = ("" if k == 0 else " · ") + seg
                _fmt_run(run, 10, bold, color)
            self.y += h + 0.04

    # ── 등급 배지(색 배지 도형) ──
    def _grade(self, block) -> None:
        gs = T.grade_style(block.grade)
        fg = _rgb(gs["fg"])
        bg = _rgb(gs["bg"])
        text = (f"{block.label}  " if block.label else "") + gs["label"]
        w = min(_CONTENT_W, max(2.4, 0.16 * len(text) + 1.0))
        hh = 0.55
        self._ensure(hh + 0.05)
        shp = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(_SAFE), Inches(self.y), Inches(w), Inches(hh))
        _fill(shp, bg)
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _p(tf.paragraphs[0], text, 15, fg, bold=True, align=PP_ALIGN.CENTER, space_after=0)
        self.y += hh + _GAP

    # ── 이미지(지도 PNG 등) ──
    def _image(self, block) -> None:
        try:
            # 너비 결정(mm 지정 있으면 인치로, 없으면 본문폭의 절반 정도)
            width_in = (block.max_width_mm / 25.4) if block.max_width_mm else 6.5
            width_in = min(width_in, _CONTENT_W)
            # 실제 비율을 알려고 임시 삽입 → 높이 계산 후 제거·재삽입
            probe = self.slide.shapes.add_picture(io.BytesIO(block.png), Inches(0), Inches(0), width=Inches(width_in))
            aspect = (probe.height / probe.width) if probe.width else 0.66
            probe._element.getparent().remove(probe._element)
            h_in = width_in * aspect
            max_h = _CONTENT_BOTTOM - _CONTENT_TOP
            if h_in > max_h:                # 너무 크면 높이에 맞춰 축소
                h_in = max_h
                width_in = h_in / aspect if aspect else width_in
            self._ensure(h_in + 0.1)
            self.slide.shapes.add_picture(
                io.BytesIO(block.png), Inches(_SAFE), Inches(self.y), width=Inches(width_in))
            self.y += h_in + _GAP
        except Exception:  # noqa: BLE001  # 이미지 깨져도 문서는 계속
            self._caption("(이미지 없음)")
        if block.caption:
            self._caption(block.caption)

    # ── 최종 면책(문서 끝) ──
    def _final_disclaimer(self) -> None:
        text = self.disclaimer_text
        h = _est_text_h(text, 100, 9) + 0.1
        # 본문 슬라이드가 하나도 없거나 공간 부족이면 면책 전용 슬라이드
        if self.slide is None or self.y + h > _CONTENT_BOTTOM:
            self._new_slide("면책 고지")
            self.section_title = "면책 고지"
        self._textbox(_SAFE, self.y, _CONTENT_W, h, text, 9, MUTED)
        self.y += h


# ── 최상위 진입점 ──────────────────────────────────────────────────────
def render_pptx(model: ReportModel) -> bytes:
    """정본 ReportModel → PPTX bytes.

    순서(pdf_renderer 미러): 표지 → (핵심요약) → 섹션들 → 면책.
    """
    prs = Presentation()
    prs.slide_width = Inches(_W)
    prs.slide_height = Inches(_H)

    deck = _Deck(prs, model.meta, model.disclaimer or T.DISCLAIMER_TEXT)

    # 표지
    deck._cover()

    # 핵심 요약(두괄식) — 있으면 먼저
    if model.exec_summary:
        deck._section(model.exec_summary)

    # 섹션들
    for sec in model.sections:
        deck._section(sec)

    # 면책
    deck._final_disclaimer()

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
