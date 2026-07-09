"""시장조사보고서 서비스 — 주변 실거래·시세·입지·수급을 통합해 심층 보고서 생성.

데이터: MolitClient(유형별 실거래 통계) + LandInfoService(용도지역·공시지가·입지) +
AI 내러티브(get_llm, best-effort). 출력: 구조화 dict / PDF(reportlab) / PPTX(python-pptx).
"""

import io
import json
import re
from datetime import datetime
from typing import Any

import structlog

from app.services.common.pdf_escape import esc as _esc

logger = structlog.get_logger(__name__)

# 주소에서 시/군/구 토큰 추출. KOSIS·SGIS는 통합시(수원·성남·고양 등)를 다르게 표기한다:
#   KOSIS 인구이동/소득 표 → 시 단위 집계('수원시', 자치구 없음)
#   SGIS stage 목록        → 시+구('수원시 장안구')
# 광역/특별시 자치구는 양쪽 모두 '강남구', 일반 시/군은 '파주시'/'○○군' 그대로.
_SIGUNGU_RE = re.compile(r"([가-힣]{1,6}(?:시|군|구))")


def _extract_region_keys(address: str | None) -> tuple[str | None, str | None]:
    """주소 → (KOSIS 시군구명, SGIS 시군구명). provider별 표기 차이를 흡수한다."""
    if not address:
        return (None, None)
    toks = _SIGUNGU_RE.findall(address)
    if not toks:
        return (None, None)
    # 시도(광역/특별/특별자치시)는 제외 — 실제 시군구만.
    si = next((t for t in reversed(toks)
               if t.endswith("시") and not t.endswith(("광역시", "특별시", "특별자치시"))), None)
    gu = next((t for t in reversed(toks) if t.endswith("구")), None)
    gun = next((t for t in reversed(toks) if t.endswith("군")), None)
    if si and gu:        # 통합시 자치구(수원시 장안구): KOSIS=시, SGIS="시 구"
        return (si, f"{si} {gu}")
    if gu:               # 광역/특별시 자치구(강남구): 양쪽 동일
        return (gu, gu)
    if si:               # 일반시(파주시)
        return (si, si)
    if gun:              # 군
        return (gun, gun)
    return (None, None)

_TRADE = [("apt", "아파트"), ("villa", "연립·다세대"), ("officetel", "오피스텔"), ("house", "단독·다가구")]
_RENT = [("apt", "아파트"), ("villa", "연립·다세대"), ("officetel", "오피스텔")]

# 면책 고지 — 모든 분석 산출물(보고서) 공통
DISCLAIMER_TEXT = (
    "본 분석결과는 참고용이며, 오류가 있을 수 있습니다. "
    "이와 관련해 사통팔땅은 어떠한 책임도 지지 않습니다. "
    "최종판단은 사용자가 최종 결정하는 것입니다."
)


PYEONG_SQM = 3.305785  # 1평 = 3.305785㎡


def _stat(values: list[float]) -> dict[str, Any]:
    vals = [v for v in values if v and v > 0]
    if not vals:
        return {"count": 0, "avg": 0, "min": 0, "max": 0}
    return {"count": len(vals), "avg": round(sum(vals) / len(vals)), "min": min(vals), "max": max(vals)}


def _per_pyeong_stat(rows: list) -> dict[str, Any]:
    """거래 행에서 평당 단가(만원/평) 통계. price_10k_won(만원) / (area_m2/3.305785)."""
    vals: list[float] = []
    for x in rows:
        p = float(x.get("price_10k_won") or 0)
        a = float(x.get("area_m2") or 0)
        if p > 0 and a > 0:
            vals.append(p / (a / PYEONG_SQM))
    s = _stat(vals)
    # 평당가는 만원 단위 정수로 반올림
    return {"count": s["count"], "avg": round(s["avg"]), "min": round(s["min"]), "max": round(s["max"])}


def _eok(man: float) -> str:
    if not man:
        return "-"
    # 손실 사업(순이익·NPV 음수)도 부호를 보존해 일관 표기("-5.0억"). 절대값 기준으로 억/만 분기.
    sign = "-" if man < 0 else ""
    a = abs(man)
    if a >= 10000:
        return f"{sign}{a / 10000:.1f}억"
    return f"{sign}{int(a):,}만"


# ── raw_data 빌더(순수 함수·네트워크 없음) ──────────────────────────────────
# 프론트(P3)·export(P4)가 dict 를 재가공하지 않도록 표(row 배열)로 평탄화한다.
# 가짜 데이터 금지: provider 가 안 준 값은 None, 미선택 분류는 키 자체를 생략한다(아래 규칙).

def _re_per_pyeong(price_10k: float, area_m2: float) -> float | None:
    """총액(만원)·면적(㎡) → 평당가(만원/평). 값이 유효할 때만 반올림 정수."""
    if price_10k and area_m2 and price_10k > 0 and area_m2 > 0:
        return round(price_10k / (area_m2 / PYEONG_SQM))
    return None


def _build_trade_table(trade: dict[str, Any]) -> list[dict[str, Any]]:
    """stats['trade'] → 행 배열. 각 유형의 건수·총액통계·평균면적·평당가를 평탄화."""
    rows: list[dict[str, Any]] = []
    for label, s in (trade or {}).items():
        pp = (s.get("per_pyeong") or {}).get("avg")
        rows.append({
            "type": label,
            "count": s.get("count", 0),
            "avg_10k": s.get("avg", 0),
            "min_10k": s.get("min", 0),
            "max_10k": s.get("max", 0),
            "avg_area_m2": s.get("avg_area_m2", 0),
            "per_pyeong_manwon": pp if pp else None,
        })
    return rows


def _build_rent_table(rent: dict[str, Any]) -> list[dict[str, Any]]:
    """stats['rent'] → 행 배열. 보증금 통계(있는 필드만)."""
    rows: list[dict[str, Any]] = []
    for label, s in (rent or {}).items():
        rows.append({
            "type": label,
            "count": s.get("count", 0),
            "avg_10k": s.get("avg", 0),
            "min_10k": s.get("min", 0),
            "max_10k": s.get("max", 0),
        })
    return rows


def _build_trend_series(apt_trend: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """apt_trend(과거→현재 정렬됨) → {ym, per_pyeong_manwon, mom_pct} 행 배열.

    mom_pct(전월대비 증감률 %) = ((현재-직전)/직전*100) 반올림 1자리. 첫 항목은 None(직전 없음).
    직전 평당가가 0/None 이면 분모 0 방지로 None.
    """
    out: list[dict[str, Any]] = []
    prev: float | None = None
    for t in (apt_trend or []):
        cur = t.get("avg_per_pyeong")
        mom = None
        if prev and cur and prev > 0:
            mom = round((cur - prev) / prev * 100, 1)
        out.append({
            "ym": t.get("ym"),
            "per_pyeong_manwon": cur if cur else None,
            "mom_pct": mom,
        })
        if cur:
            prev = cur
    return out


def _build_population_block(demographics: dict[str, Any] | None) -> dict[str, Any] | None:
    """demographics.population/migration → 표 평탄화 블록.

    미선택 분류(population 비어 있고 data_source 없음)면 None 반환 → 호출측에서 키 생략.
    선택했으나 provider 실패면 data_source='unavailable'/'fallback' 등으로 정직 표기.
    """
    pop = (demographics or {}).get("population") or {}
    mig = (demographics or {}).get("migration") or {}
    pop_src = pop.get("data_source")
    mig_src = mig.get("data_source")
    # 둘 다 data_source 가 없으면 = 미선택(호출 자체를 안 함) → 키 생략 신호로 None.
    if pop_src is None and mig_src is None:
        return None

    # 연령분포: SGIS 는 {label: count} dict(예: {"0-9": n, ...}) → 행 배열.
    age_dist = pop.get("age_distribution") or {}
    age_rows = [{"label": k, "count": v} for k, v in age_dist.items()] if isinstance(age_dist, dict) else []

    # 가구원수 분포: {"1_person","2_person","3_person","4_over"} 비율(%) → 행 배열.
    #   SGIS 미제공이라 평균 가구원수 기반 '추정'(estimated=true) — note 로도 명시됨.
    ht = pop.get("household_types") or {}
    _ht_label = {"1_person": "1인", "2_person": "2인", "3_person": "3인", "4_over": "4인+"}
    household_rows = (
        [{"label": _ht_label.get(k, k), "ratio": v, "estimated": True} for k, v in ht.items()]
        if isinstance(ht, dict) else []
    )

    return {
        "summary": {
            "total_population": pop.get("total_population") or None,
            "household_count": pop.get("household_count") or None,
            "avg_household_size": pop.get("avg_household_size") or None,
        },
        "age_distribution": age_rows,
        "household_types": household_rows,
        "migration": {
            "total_inflow": mig.get("total_inflow") if mig_src is not None else None,
            "total_outflow": mig.get("total_outflow") if mig_src is not None else None,
            "net_migration": mig.get("net_migration") if mig_src is not None else None,
        },
        "source": "통계청 인구주택총조사 / 국내인구이동통계",
        # population 자체의 출처. 미선택(호출 안 함)이라 None 이면 정직하게 unavailable.
        "data_source": pop_src or "unavailable",
        "migration_data_source": mig_src or "unavailable",
    }


def _build_income_block(demographics: dict[str, Any] | None) -> dict[str, Any] | None:
    """demographics.macro_income → 소득 블록. 미선택(data_source 없음)이면 None(키 생략).

    KOSIS 국세청 표는 평균 총급여만 산출 가능(중위는 평균×0.85 추정). 인원/총급여 원수치는
    모델에 보존되지 않으므로 basis.persons/total_salary_10k 는 None(가짜 금지). bracket_ratio 미제공→null.
    """
    mi = (demographics or {}).get("macro_income") or {}
    src = mi.get("data_source")
    if src is None:
        return None
    avg = mi.get("avg_income_10k") or None
    med = mi.get("median_income_10k") or None
    # income_bracket_ratio 가 비어 있으면(미제공) 정직 null.
    bracket = mi.get("income_bracket_ratio") or None
    return {
        "avg_income_10k": avg,
        "median_income_10k": med,
        # 중위는 평균×0.85 결정론 추정값 — 실측 아님을 명시.
        "median_estimated": True,
        # 인원/총급여 원수치는 모델에 보존되지 않음(KosisClient 내부 계산값) → 정직 None.
        "basis": {"persons": None, "total_salary_10k": None},
        "bracket_ratio": bracket,
        "source": "국세청 근로소득",
        "data_source": src,
    }


# ── 타겟 프로파일·Executive Summary 공용 헬퍼(순수 함수·네트워크 없음) ─────────────
# K-Atlas(micro_finance) 대체: 이미 라이브인 실데이터(SGIS 연령/가구·KOSIS 소득·SEMAS 상권·
# 입지 POI)로 '타겟 프로파일'을 조립한다. 각 축은 data_source(live/fallback/unavailable)로
# 정직 표기하고, 신용·카드소비 등 마이크로 금융은 'PREMIUM 제휴 예정'으로 정직 강등한다.

def _has_source(ds: Any) -> bool:
    """data_source 가 실데이터(비-미확보)인지. '', None, unavailable, mock 은 미확보로 본다.

    ★무목업: 'mock'(KOSIS/SGIS 무키 시 반환되는 테스트용 가짜값)을 실데이터로 취급하면
    타겟 프로파일·시니어 근거에 가짜 수치(예: 소득 4,620만원)가 '실데이터'로 노출된다.
    'fallback'(전국 평균 근사·공개 고지)은 유지하되 'mock'은 미확보로 차단한다.
    """
    return str(ds or "").lower() not in ("", "unavailable", "mock")


def _income_tier_label(avg_income_10k: float | None) -> str | None:
    """연평균 소득(만원) → 소득 수준 라벨(전국 대비 참고 구간·추정). 공식 소득분위 아님(정직).

    통계청 공식 10분위 컷오프는 연도·표본별로 달라 데이터 근거 없이 확정하지 않는다.
    대신 널리 알려진 가구 연소득 대역으로 rough 구간만 표기하고 '추정'을 명시한다.
    """
    if not avg_income_10k or avg_income_10k <= 0:
        return None
    v = float(avg_income_10k)
    if v < 3000:
        return "중하위 소득권(추정)"
    if v < 4500:
        return "중위 소득권(추정)"
    if v < 6000:
        return "중상위 소득권(추정)"
    return "상위 소득권(추정)"


def _roi_grade(roi_percent: float | None) -> tuple[str, str]:
    """개략 ROI(%) → (등급, 라벨). 근거(수지) 없으면 ('-', '데이터 없음')."""
    if roi_percent is None:
        return ("-", "데이터 없음")
    r = float(roi_percent)
    if r >= 20:
        return ("A", "우수")
    if r >= 10:
        return ("B", "양호")
    if r >= 0:
        return ("C", "보통")
    return ("D", "주의")


def _market_strength(net_migration: int | None, trend_series: list[dict[str, Any]] | None) -> tuple[str, str]:
    """시장 강도 = 순유입(수요 유입) + 최근 시세 방향성. (등급, 라벨). 근거 없으면 ('-', '데이터 없음')."""
    # 최근 추이 방향(마지막 mom_pct 부호)
    last_mom = None
    for t in reversed(trend_series or []):
        if t.get("mom_pct") is not None:
            last_mom = t.get("mom_pct")
            break
    if net_migration is None and last_mom is None:
        return ("-", "데이터 없음")
    score = 0
    if net_migration is not None:
        score += 1 if net_migration > 0 else (-1 if net_migration < 0 else 0)
    if last_mom is not None:
        score += 1 if last_mom > 0 else (-1 if last_mom < 0 else 0)
    if score >= 1:
        return ("강세", "수요 유입·시세 상승 우위")
    if score <= -1:
        return ("약세", "수요 유출·시세 조정 우위")
    return ("보합", "수요·시세 혼조(중립)")


def _investment_opinion(roi_percent: float | None, afford_verdict: str | None) -> str:
    """결정론 투자의견(Go/보류/재검토). 개략 ROI·지불여력 판정 기반(전문검토 전제)."""
    if roi_percent is None:
        return "재검토 — 사업성 데이터 부족(개략 수지 산출 불가)"
    r = float(roi_percent)
    over = afford_verdict == "over_band"
    if r >= 10 and not over:
        return "Go(추진 검토) — 개략 수익성 양호"
    if r >= 0:
        base = "보류(조건부) — 수익성 보통, 원가·분양가 가정 민감도 점검 필요"
        return base + (" · 지불여력 초과(미분양 위험) 주의" if over else "")
    return "재검토 — 개략 수익성 음수(총사업비 대비 분양수익 부족)"


def _build_target_profile(
    demographics: dict[str, Any] | None,
    commercial: dict[str, Any] | None,
    commercial_src: str,
    infra: dict[str, Any] | None,
) -> dict[str, Any]:
    """실데이터 5축 타겟 프로파일 조립. 각 축 data_source 정직 표기(무목업).

    ① 주력연령(SGIS age_distribution 최빈대) ② 주력가구(household_types 최빈)
    ③ 소득분위(macro_income avg/median→라벨) ④ 상권특성(SEMAS analyze_commercial_area)
    ⑤ 입지(POI/subway 요약). 신용·카드소비(K-Atlas)는 PREMIUM 제휴 예정으로 강등.
    """
    dm = demographics or {}
    pop = dm.get("population") or {}
    mi = dm.get("macro_income") or {}
    pop_src = pop.get("data_source")
    inc_src = mi.get("data_source")
    tp: dict[str, Any] = {}

    # ① 주력 연령대 — 연령분포 최빈 밴드
    #   value/detail = 화면·타일 표시용(프론트 타겟프로파일 카드 계약: {label,value,detail}).
    #   band/count = 렌더러(PDF/PPTX/DOCX)용 원시필드(보존). 둘 다 유지(무회귀).
    age = pop.get("age_distribution") or {}
    if isinstance(age, dict) and age and _has_source(pop_src):
        try:
            top = max(((k, float(v or 0)) for k, v in age.items()), key=lambda kv: kv[1])
            _band = str(top[0])
            _m = re.match(r"^(\d+)0\s*[-~]", _band)          # '30-39' → '30대'
            _band_label = f"{_m.group(1)}0대" if _m else _band
            tp["primary_age"] = {
                "band": top[0], "count": int(top[1]), "data_source": pop_src,
                "value": _band_label, "detail": f"{int(top[1]):,}명",
            }
        except (ValueError, TypeError):
            tp["primary_age"] = {"data_source": "unavailable"}
    else:
        tp["primary_age"] = {"data_source": "unavailable"}

    # ② 주력 가구 유형 — 가구원수 분포 최빈(SGIS 미제공→평균 가구원수 기반 추정치)
    ht = pop.get("household_types") or {}
    _ht_label = {"1_person": "1인 가구", "2_person": "2인 가구", "3_person": "3인 가구", "4_over": "4인+ 가구"}
    if isinstance(ht, dict) and ht and _has_source(pop_src):
        try:
            top = max(((k, float(v or 0)) for k, v in ht.items()), key=lambda kv: kv[1])
            _ht_type = _ht_label.get(top[0], top[0])
            # ratio가 0~1(비율)이면 %로, 이미 %면 그대로(방어적).
            _r = top[1] * 100 if 0 < top[1] <= 1 else top[1]
            tp["primary_household"] = {
                "type": _ht_type, "ratio": top[1],
                "estimated": True, "data_source": pop_src,
                "value": _ht_type, "detail": f"{_r:.0f}% (추정)" if _r else "추정",
            }
        except (ValueError, TypeError):
            tp["primary_household"] = {"data_source": "unavailable"}
    else:
        tp["primary_household"] = {"data_source": "unavailable"}

    # ③ 소득 분위(라벨) — 거시 평균/중위 소득 기반 참고 구간(추정·공식 분위 아님)
    avg = mi.get("avg_income_10k")
    med = mi.get("median_income_10k")
    if avg and _has_source(inc_src):
        _tier = _income_tier_label(avg)
        tp["income_tier"] = {
            "avg_income_10k": avg,
            "median_income_10k": med or None,
            "tier_label": _tier,
            "basis": "전국 가구소득 대비 참고 구간(추정) — 통계청 공식 소득분위 아님",
            "data_source": inc_src,
            "value": _tier or "—", "detail": f"연 {int(avg):,}만원" if avg else None,
        }
    else:
        tp["income_tier"] = {"data_source": "unavailable"}

    # ④ 상권 특성 — SEMAS analyze_commercial_area 결과(라이브만 채택, 실패/무키는 정직 강등)
    if commercial and str(commercial_src) == "live":
        _stores = commercial.get("total_stores")
        _grade = commercial.get("grade")
        tp["commercial"] = {
            "total_stores": _stores,
            "category_distribution": commercial.get("category_distribution"),
            "vitality_score": commercial.get("vitality_score"),
            "grade": _grade,
            "data_source": "live",
            "value": str(_grade) if _grade else "상권 형성",
            "detail": f"점포 {int(_stores):,}개" if _stores else None,
        }
    else:
        tp["commercial"] = {
            "data_source": "unavailable",
            "note": "상권(SEMAS) 데이터 미확보 — 키 미설정/좌표 부재/미확보(정직 표기)",
        }

    # ⑤ 입지 — 기존 인프라(지하철·학교 POI) 요약
    inf = infra if isinstance(infra, dict) else {}
    subway = inf.get("nearest_subway") or {}
    schools = inf.get("schools") or []
    has_loc = bool(subway.get("name") or schools)
    _sub = subway.get("name")
    _sch = len(schools) if isinstance(schools, list) else 0
    # value = 입지 핵심(역세권/학군 요약), detail = 지하철명/학교수(정직).
    _loc_bits = []
    if _sub:
        _loc_bits.append("역세권")
    if _sch:
        _loc_bits.append("학군")
    tp["location"] = {
        "nearest_subway": _sub,
        "subway_distance_m": subway.get("distance_m"),
        "school_count": _sch if isinstance(schools, list) else None,
        "data_source": "live" if has_loc else "unavailable",
        "value": "·".join(_loc_bits) if _loc_bits else ("입지 정보" if has_loc else None),
        "detail": (f"{_sub}" + (f"·학교 {_sch}" if _sch else "")) if _sub else (f"학교 {_sch}" if _sch else None),
    }

    # 신용·카드소비 등 마이크로 금융(K-Atlas): PREMIUM 제휴 예정으로 정직 강등(가짜값 금지)
    tp["premium"] = {
        "credit_score": {"status": "PREMIUM 제휴 예정", "data_source": "unavailable"},
        "card_spending": {"status": "PREMIUM 제휴 예정", "data_source": "unavailable"},
        "note": "신용평점·카드소비 등 마이크로 금융데이터(K-Atlas 등)는 PREMIUM 제휴 연동 예정 — 현재 미제공(정직 표기).",
    }
    return tp


def _exec_summary(rep: dict[str, Any]) -> dict[str, Any]:
    """두괄식 Executive Summary용 핵심 KPI·등급·투자의견(순수 함수·재계산 최소).

    핵심 KPI: 아파트 평당시세·순유입·적정분양가·ROI. 등급: 시장강도·사업성.
    투자의견: 결정론(Go/보류/재검토) + senior_insight(있으면) 1문장.
    """
    # 아파트 평당 시세(만원/평)
    apt_pp = None
    _apt = (rep.get("trade") or {}).get("아파트")
    if isinstance(_apt, dict):
        apt_pp = (_apt.get("per_pyeong") or {}).get("avg")

    # 순유입(전입-전출) — ★미확보('unavailable'/'mock')는 net_migration이 0으로 채워져 오므로,
    #   출처(migration_data_source)를 함께 게이트해야 가짜 '0명'·근거없는 시장강도 등급을 막는다(무목업).
    net_mig = None
    pop = (rep.get("raw_data") or {}).get("population") or {}
    mig = pop.get("migration") or {}
    if _has_source(pop.get("migration_data_source")) and mig.get("net_migration") is not None:
        net_mig = mig.get("net_migration")

    # 적정 분양가(거래사례비교) — 미확보면 None
    pb = rep.get("pricing_band") or {}
    fair = pb.get("fair_price_10k") if _has_source(pb.get("data_source")) else None
    afford_verdict = pb.get("affordability_verdict")

    # 개략 ROI(%)
    fe = rep.get("feasibility_analysis") or {}
    roi = None
    if isinstance(fe, dict) and not fe.get("error"):
        roi = (fe.get("financials") or {}).get("roi_percent")

    # 시세 추이(방향성) — raw_data trend_series
    trend = ((rep.get("raw_data") or {}).get("real_estate") or {}).get("trend_series")

    biz_grade, biz_label = _roi_grade(roi)
    mkt_grade, mkt_label = _market_strength(net_mig, trend)
    opinion = _investment_opinion(roi, afford_verdict)
    # senior_insight 투자의견 1문장(있으면 병기 — 전문가 서술)
    si = rep.get("senior_insight") or {}
    si_opinion = si.get("investment_insight") or si.get("timing_recommendation")

    return {
        "market_strength": {"grade": mkt_grade, "label": mkt_label},
        "business_grade": {"grade": biz_grade, "label": biz_label},
        "kpi": {
            "apt_per_pyeong_manwon": apt_pp,
            "net_migration": net_mig,
            "fair_price_10k": fair,
            "roi_percent": roi,
        },
        "opinion": opinion,
        "expert_opinion": (str(si_opinion) if si_opinion else None),
    }


def _insight_text(rep: dict[str, Any], senior_key: str, narrative_key: str | None = None) -> str | None:
    """섹션 내러티브 1~2문장 선택: senior_insight(전용 인터프리터) 우선, 없으면 인라인 narrative 폴백.

    둘 다 없으면 None(호출측이 'AI 분석 미포함' 정직 표기). 표만 있던 섹션의 수치→해석 주입용.
    """
    si = rep.get("senior_insight") or {}
    val = si.get(senior_key)
    if val and str(val).strip():
        return str(val).strip()
    if narrative_key:
        nar = rep.get("narrative") or {}
        nv = nar.get(narrative_key)
        if nv and str(nv).strip():
            return str(nv).strip()
    return None


class MarketReportService:
    def __init__(self) -> None:
        from apps.api.integrations.molit_client import MolitClient

        self.molit = MolitClient()

    def _months(self, n: int = 3) -> list[str]:
        now = datetime.now()
        y, m = now.year, now.month - 1  # 현재월 신고지연 → 직전월부터
        if m == 0:
            m = 12
            y -= 1
        out = []
        for _ in range(n):
            out.append(f"{y}{m:02d}")
            m -= 1
            if m == 0:
                m = 12
                y -= 1
        return out

    async def _category_stats(self, lawd_cd: str) -> dict[str, Any]:
        import asyncio

        # ★MOLIT LAWD_CD 는 시군구 5자리. 10자리 법정동/bcode 가 들어오면 [:5]로 정규화하지 않으면
        #   실거래가 전부 빈 결과가 되어 trade·comparable_trade(적정분양가 앵커)가 통째로 누락된다.
        lawd_cd = (lawd_cd or "")[:5]
        months = self._months(3)
        trade: dict[str, Any] = {}
        rent: dict[str, Any] = {}

        async def trade_one(pt: str, label: str):
            rows: list = []
            res = await asyncio.gather(*[self.molit.get_transactions(lawd_cd, ym, prop_type=pt, num_rows=1000) for ym in months], return_exceptions=True)
            for r in res:
                if isinstance(r, list):
                    rows.extend(r)
            prices = [float(x.get("price_10k_won") or 0) for x in rows]
            areas = [float(x.get("area_m2") or 0) for x in rows]
            return label, {
                **_stat(prices),
                "avg_area_m2": round(sum(a for a in areas if a > 0) / max(1, len([a for a in areas if a > 0])), 1) if areas else 0,
                "per_pyeong": _per_pyeong_stat(rows),  # 평당가(만원/평) — 면적 정규화 시세
            }

        async def rent_one(pt: str, label: str):
            rows: list = []
            res = await asyncio.gather(*[self.molit.get_rent_transactions(lawd_cd, ym, prop_type=pt, num_rows=1000) for ym in months], return_exceptions=True)
            for r in res:
                if isinstance(r, list):
                    rows.extend(r)
            dep = [float(x.get("deposit_10k_won") or 0) for x in rows]
            return label, {**_stat(dep), "count": len([d for d in dep if d > 0])}

        # 아파트 매매 월별 추이(시세 추이 차트용)
        async def apt_month(ym: str):
            try:
                rows = await self.molit.get_transactions(lawd_cd, ym, prop_type="apt", num_rows=1000)
            except Exception:  # noqa: BLE001
                rows = []
            prices = [float(x.get("price_10k_won") or 0) for x in rows if (x.get("price_10k_won") or 0) > 0]
            pp = _per_pyeong_stat(rows)
            return {
                "ym": ym,
                "avg": round(sum(prices) / len(prices)) if prices else 0,  # 총액 평균(만원)
                "avg_per_pyeong": pp["avg"],  # 평당가(만원/평) — 추이 기준
                "count": len(prices),
            }

        tr = await asyncio.gather(*[trade_one(pt, lb) for pt, lb in _TRADE])
        rr = await asyncio.gather(*[rent_one(pt, lb) for pt, lb in _RENT])
        trend = await asyncio.gather(*[apt_month(ym) for ym in months])
        trade = dict(tr)
        rent = dict(rr)
        # 추이는 과거→현재 순으로
        trend_sorted = sorted(trend, key=lambda t: t["ym"])
        return {"months": months, "trade": trade, "rent": rent, "apt_trend": trend_sorted}

    async def _narrative(self, ctx: dict[str, Any]) -> dict[str, Any]:
        """AI 시장 해석(요약·기회·리스크). 실패 시 구조화 폴백."""
        try:
            from langchain_core.messages import HumanMessage, SystemMessage

            from app.services.ai.llm_provider import get_llm

            llm = get_llm(timeout=40, max_tokens=1500)
            # 근거 그라운딩: 대상지 주소 기반 지역 시세 벤치마크를 근거로 주입한다.
            # 전용 인터프리터(MarketInterpreter)와 동일 근거원인 BaseInterpreter._regional_benchmark
            # (정적·sync·키불필요·결정적 로컬 테이블 — 블로킹 I/O 없음)를 재사용해, 인라인 경로도
            # 전용 경로와 같은 근거를 보게 한다. 실패는 내러티브 무손상(best-effort).
            benchmark = None
            try:
                from app.services.ai.base_interpreter import BaseInterpreter
                benchmark = BaseInterpreter._regional_benchmark(address=str(ctx.get("address", "")))
            except Exception:  # noqa: BLE001 — 벤치마크 조회 실패는 내러티브 무손상
                benchmark = None
            sys = ("당신은 감정평가사·분양대행 실무 관점을 겸비한 부동산 개발·시장분석 전문가다. 제공된 실거래·시세·입지 데이터와 인구 이동, 연령대, 평균 소득 데이터를 종합하여 "
                   "한국어 JSON으로 답하라. 키: summary(시장요약 3~4문장), opportunities(기회 2~3개 배열), "
                   "risks(리스크 2~3개 배열), price_trend(가격동향 2문장), target_persona(추천 분양 타겟 고객층 2문장). "
                   "★모든 거래시세·분양가는 반드시 평당가(만원/평) 기준으로 서술하라. 총액(억원)이 아닌 "
                   "평당 단가를 사용한다. 예: '아파트 평당 약 1,800만원'. 데이터 단위는 만원/평이다. "
                   "★target_persona에는 유입 인구의 주 연령대, 거시적 평균 소득을 고려해 가장 분양 가능성이 높은 고객의 직업군/가구형태/특화설계 제안을 포함하라. "
                   "★'지역 시세 벤치마크(참고 근거)'가 제공되면 데이터 적정성 판단에 참고하되, 그 값을 '지역 평균 분양가' 같은 확정 사실로 단정하지 말고 부득이 언급 시 '참고 추정'임을 명시하라. 수집 데이터가 없는 항목은 단정하지 말라(무근거 추정 금지).")
            usr = f"## 시장 데이터\n{json.dumps(ctx, ensure_ascii=False)[:4000]}"
            if benchmark:
                usr += f"\n\n## 지역 시세 벤치마크(참고 근거)\n{benchmark}"
            resp = await llm.ainvoke([SystemMessage(content=sys), HumanMessage(content=usr)])
            # 계측: BaseInterpreter 밖 직접 호출도 동일하게 토큰·과금 기록(best-effort)
            from app.services.ai.base_interpreter import record_llm_response_billing
            await record_llm_response_billing(llm, resp, service="market_report")
            raw = resp.content if hasattr(resp, "content") else str(resp)
            txt = raw.strip()
            if txt.startswith("```"):
                txt = txt.split("```")[1].lstrip("json").strip() if "```" in txt[3:] else txt.strip("`")
            data = json.loads(txt)
            return data
        except Exception as e:  # noqa: BLE001
            logger.warning("시장 내러티브 생성 실패, 구조화 폴백", err=str(e)[:80])
            return {"summary": "수집된 실거래·시세 데이터를 기반으로 한 시장 현황입니다.", "opportunities": [], "risks": [], "price_trend": "", "target_persona": "데이터 기반 타겟팅 분석 불가"}

    async def _nearby_presale_84_price(
        self, lawd_cd: str, coords: Any,
    ) -> tuple[float | None, str]:
        """주변 신규 분양가(청약홈) → 84㎡급 대표 분양총액(만원) 중앙값. (값, 출처).

        거래사례비교 1차 보강. PresaleService.nearby(거리필터 공고) → 가까운 아파트 공고 상위 3개의
        detail(주택형별 분양총액 price_man·공급면적)에서 84㎡급(공급 100~125㎡) 분양가를 모아 중앙값.
        키 미설정/데이터 없음/타임아웃이면 (None, 'unavailable') 정직 반환(가짜값 금지).
        보고서 지연 방지를 위해 하드타임아웃 적용.
        """
        import asyncio as _aio
        try:
            from app.services.land_intelligence.presale_service import PresaleService, area_from_lawd
        except Exception:  # noqa: BLE001
            return None, "unavailable"
        _lat = coords.get("lat") if isinstance(coords, dict) else None
        _lon = (coords.get("lon") or coords.get("lng")) if isinstance(coords, dict) else None
        try:
            svc = PresaleService()
            near = await _aio.wait_for(
                svc.nearby(_lat, _lon, area_from_lawd(lawd_cd),
                           radius_m=3000, months_back=12, max_markers=8),
                timeout=12.0)
            if not near.get("available"):
                return None, "unavailable"
            picks = [it for it in (near.get("items") or []) if it.get("house_manage_no")][:3]
            if not picks:
                return None, "unavailable"
            details = await _aio.wait_for(_aio.gather(*[
                svc.detail(it.get("house_manage_no", ""), it.get("pblanc_no", ""), it.get("product", "apt"))
                for it in picks], return_exceptions=True), timeout=15.0)
            cands: list[float] = []
            for d in details:
                if not isinstance(d, dict) or not d.get("available"):
                    continue
                for m in d.get("models", []):
                    try:
                        amt = float(m.get("price_man"))
                        ar = float(m.get("supply_area_m2"))
                    except (TypeError, ValueError):
                        continue
                    # 전용 84㎡급 ≈ 공급 100~125㎡ — 실거래 84㎡ 기준가와 일관되게 비교.
                    if amt > 0 and 100.0 <= ar <= 125.0:
                        cands.append(amt)
            if not cands:
                return None, "unavailable"
            cands.sort()
            return float(cands[len(cands) // 2]), "live"  # 중앙값
        except Exception:  # noqa: BLE001 — 타임아웃·네트워크·파싱 실패 시 정직 None
            return None, "unavailable"

    async def _commercial_area(self, coords: Any) -> tuple[dict[str, Any] | None, str]:
        """상권 활성도(SEMAS) — CommercialAreaService.analyze_commercial_area(lat,lon,500).

        ★신규 배선(과제1 ④): land_info_service.py 의 인스턴스화 패턴(self.commercial=
        CommercialAreaService())을 그대로 재사용한다. 좌표 부재/키 미설정/미확보/타임아웃은
        (None, 'unavailable') 정직 반환(가짜값 금지). 보고서 지연 방지 하드타임아웃.
        """
        import asyncio as _aio
        _lat = coords.get("lat") if isinstance(coords, dict) else None
        _lon = (coords.get("lon") or coords.get("lng")) if isinstance(coords, dict) else None
        if not (_lat and _lon):
            return None, "unavailable"
        try:
            from app.services.external_api.commercial_area_service import (
                CommercialAreaService,
            )
            res = await _aio.wait_for(
                CommercialAreaService().analyze_commercial_area(float(_lat), float(_lon), 500),
                timeout=20.0)
            if not res or not isinstance(res, dict):
                return None, "unavailable"
            return res, "live"
        except Exception:  # noqa: BLE001 — 키/네트워크/타임아웃 실패는 정직 미확보
            return None, "unavailable"

    @staticmethod
    def _senior_prior_context(
        demographics: dict[str, Any] | None,
        feasibility: dict[str, Any] | None,
        pricing_band: dict[str, Any] | None,
        unit_mix: dict[str, Any] | None,
    ) -> str | None:
        """MarketInterpreter 근거블록(prior_context) — 인구·수급·사업성·가격밴드·수요MD 요약.

        전용 인터프리터의 _extract_compact_data 가 읽지 않는 축(인구이동·소득·수지·평형MD)을
        수집·검증 근거로 프롬프트에 주입한다(가짜 금지·실값만). 근거 없으면 None.
        """
        lines: list[str] = []
        dm = demographics or {}
        pop = dm.get("population") or {}
        mig = dm.get("migration") or {}
        mi = dm.get("macro_income") or {}
        if _has_source(mig.get("data_source")) and mig.get("net_migration") is not None:
            lines.append(
                f"순유입(전입-전출): {mig.get('net_migration'):,}명 "
                f"(전입 {mig.get('total_inflow')}, 전출 {mig.get('total_outflow')})")
        age = pop.get("age_distribution") or {}
        if isinstance(age, dict) and age and _has_source(pop.get("data_source")):
            try:
                top = max(((k, float(v or 0)) for k, v in age.items()), key=lambda kv: kv[1])
                lines.append(f"주력 연령대: {top[0]} ({int(top[1]):,}명)")
            except (ValueError, TypeError):
                pass
        if mi.get("avg_income_10k") and _has_source(mi.get("data_source")):
            lines.append(f"거시 평균 연소득: {mi.get('avg_income_10k'):,}만원(국세청 근로소득)")
        fin = (feasibility or {}).get("financials") or {}
        if fin and fin.get("roi_percent") is not None and not (feasibility or {}).get("error"):
            lines.append(
                f"개략 사업성(참고 추정): ROI {fin.get('roi_percent')}% · "
                f"총사업비 {fin.get('total_cost_10k')}만원 · NPV {fin.get('npv_10k')}만원")
        pb = pricing_band or {}
        if pb.get("fair_price_10k") and _has_source(pb.get("data_source")):
            lines.append(
                f"적정 분양가(거래사례비교·84㎡급): {pb.get('fair_price_10k')}만원 · "
                f"지불여력 판정 {pb.get('affordability_verdict')}")
        um = unit_mix or {}
        if um.get("recommended_mix") and _has_source(um.get("data_source")):
            lines.append(
                f"수요기반 권장 평형 배분: {um.get('recommended_mix')} "
                f"(주력 {um.get('dominant_band')})")
        if not lines:
            return None
        return "## 인구·수급·사업성 근거(수집·검증 데이터)\n" + "\n".join(f"- {ln}" for ln in lines)

    async def _senior_market_insight(
        self, *, address: str, zone_type: str | None, land_area_sqm: float,
        stats_trade: dict[str, Any], official_price: float | None,
        pricing_band: dict[str, Any] | None, unit_mix: dict[str, Any] | None,
        demographics: dict[str, Any] | None, feasibility: dict[str, Any] | None,
    ) -> dict[str, str] | None:
        """전용 MarketInterpreter(감정평가/분양대행 페르소나)로 정밀 내러티브(6키) 산출.

        comprehensive_analysis_service.py 의 호출패턴(MarketInterpreter().generate_interpretation)
        을 그대로 재사용한다. 실거래(transaction_prices)·공시지가(land_prices)·적정분양가
        (sale_prices)·실효 용적률(effective_far)을 인터프리터가 읽는 키로 매핑하고, 인구·수지·
        수요MD는 prior_context 근거로 주입한다. 미가용/무키/예외는 None(정직 폴백).
        """
        try:
            from app.services.ai.market_interpreter import MarketInterpreter
        except Exception:  # noqa: BLE001 — 인터프리터 import 실패는 graceful 미가용
            return None
        # 실거래 통계 → 인터프리터 transaction_prices 키(avg/max/min = 만원)
        txn: dict[str, Any] = {}
        for label, s in (stats_trade or {}).items():
            if isinstance(s, dict) and s.get("count"):
                txn[label] = {
                    "count": s.get("count"),
                    "avg_price_10k": s.get("avg"),
                    "max_price_10k": s.get("max"),
                    "min_price_10k": s.get("min"),
                }
        # 공시지가(원/㎡) → land_prices
        land_prices: dict[str, Any] = {}
        if official_price and official_price > 0:
            land_prices = {
                "official_price_per_sqm": official_price,
                "official_price_per_pyeong": round(official_price * PYEONG_SQM),
            }
        # 적정 분양가(pricing_band) → sale_prices 1건(84㎡ 총액 → 평당 환산)
        sale_prices: list[dict[str, Any]] = []
        _pb = pricing_band or {}
        _fp = _pb.get("fair_price_10k")
        if _fp and _has_source(_pb.get("data_source")):
            _pp = round(_fp / (84.0 / PYEONG_SQM))  # 84㎡ 총액(만원) → 평당(만원/평)
            sale_prices.append({
                "type_name": "적정 분양가(거래사례비교·84㎡급)",
                "sale_price_per_pyeong_man": _pp,
                "sale_price_per_sqm_man": round(_fp / 84.0, 1),
            })
        # 실효 용적률/건폐율(feasibility massing) → effective_far
        _mass = (feasibility or {}).get("massing") or {}
        effective_far: dict[str, Any] = {}
        if _mass.get("estimated_far"):
            effective_far = {
                "effective_far_pct": _mass.get("estimated_far"),
                "effective_bcr_pct": _mass.get("estimated_bca"),
            }
        market_data = {
            "address": address,
            "zone_type": zone_type,
            "land_area_sqm": land_area_sqm,
            "transaction_prices": txn,
            "land_prices": land_prices,
            "sale_prices": sale_prices,
            "effective_far": effective_far,
        }
        prior = self._senior_prior_context(demographics, feasibility, pricing_band, unit_mix)
        try:
            insight = await MarketInterpreter().generate_interpretation(
                market_data, prior_context=prior)
        except Exception:  # noqa: BLE001 — 인터프리터 호출 실패는 정직 폴백(None)
            return None
        # _invoke 는 무키/파싱실패 시 {} 반환 → 비면 None(정직) 처리(렌더러가 인라인 폴백).
        return insight or None

    @staticmethod
    def _attach_senior(
        result: dict[str, Any], feasibility: dict[str, Any] | None, *,
        official_price: float | None = None, land_area_sqm: float | None = None,
        zone_type: str | None = None, real_pp: float | None = None,
    ) -> None:
        """시장 보고서에 시니어 다전문가(감정평가·도시계획·금융) 자문 첨부.

        과제2: 단일 finance → attach_senior_consultation_multi(['appraisal','urban','finance']).
        실값 주입: 감정평가=공시지가 기준 토지가액(원=공시지가/㎡×대지면적) / 금융=총사업비(원) /
        도시계획=용도·면적(정비사업 정량입력 아님→프레임워크·근거 흐름·무목업 정직).
        시장 보고서는 NOI·부채 등 대주지표 쌍이 없어 finance 정량 verdict는 보통 None(정직).
        ★무회귀: 절대 raise 안 함(graceful)."""
        try:
            from app.services.senior_agents.consultation_hook import (
                attach_senior_consultation_multi,
            )

            _sr_inputs: dict[str, Any] = {}
            # 금융: 총사업비(원) = 만원 × 10000(financial evaluator는 원 단위)
            _fin = (feasibility or {}).get("financials") if isinstance(feasibility, dict) else None
            if isinstance(_fin, dict):
                _tc = _fin.get("total_cost_10k")
                if isinstance(_tc, (int, float)) and not isinstance(_tc, bool) and _tc > 0:
                    _sr_inputs["total_cost"] = float(_tc) * 10000.0
            # 감정평가: 토지 공시지가 기준 감정가(원) = 공시지가(원/㎡) × 대지면적(㎡)
            #   (건물 없는 개발부지 → 토지만 → appraisal evaluator가 WARN·토지가액 명시. 무목업.)
            if (official_price and official_price > 0
                    and land_area_sqm and land_area_sqm > 0):
                _sr_inputs["land_appraised_total"] = float(official_price) * float(land_area_sqm)
            # 도시계획·감정 참고 실값(context) — 현 evaluator는 정비사업 비례율/종전감정 입력만
            #   정량화하므로 아래 값은 소비되지 않고 framework·근거만 흐른다. 그래도 실값을 함께
            #   실어 자문 컨텍스트를 정직하게 전달한다(향후 evaluator 확장 시 자동 활용).
            if zone_type:
                _sr_inputs["zone_type"] = zone_type
            if land_area_sqm and land_area_sqm > 0:
                _sr_inputs["land_area_sqm"] = float(land_area_sqm)
            if real_pp and real_pp > 0:
                _sr_inputs["comparable_price_per_pyeong_10k"] = float(real_pp)
            result["senior_consultation"] = attach_senior_consultation_multi(
                ["appraisal", "urban", "finance"], _sr_inputs)
        except Exception:  # noqa: BLE001 — 시니어 자문 첨부 실패는 시장 분석 무손상
            pass

    async def build_report(self, address: str, lawd_cd: str, pnu: str | None = None, use_llm: bool = True, options: dict | None = None, parcels: list[dict] | None = None) -> dict[str, Any]:
        from app.services.land_intelligence.land_info_service import LandInfoService

        options = options or {}
        use_sgis = options.get("sgis", False)
        use_kosis = options.get("kosis", False)
        # 항목 단위 게이팅: 프론트(P1)가 보내는 세부 선택. 없을 수도 있음(하위호환).
        #   detail = {pop_age, pop_household, pop_migration, income_avg, income_basis}
        #   detail 이 제공되면 항목 기준으로, 없으면 기존 분류 boolean(use_sgis/use_kosis)으로 폴백한다.
        #   → detail 미전달 시 기존 호출 동작 100% 무회귀.
        detail = options.get("detail") or {}

        def _want(key: str, fallback: bool) -> bool:
            """detail 에 항목이 명시되면 그 값, 없으면 fallback(분류 boolean)."""
            return bool(detail[key]) if key in detail else fallback

        # 인구규모/연령/가구원수는 SGIS 단일 호출이라 연령·가구 둘 중 하나라도 ON 이면 호출.
        want_pop = (
            (_want("pop_age", use_sgis) or _want("pop_household", use_sgis))
            if ("pop_age" in detail or "pop_household" in detail)
            else use_sgis
        )
        # 인구이동: detail.pop_migration 명시 시 그 값, 비면 기존 use_sgis/use_kosis 폴백.
        want_mig = _want("pop_migration", use_sgis or use_kosis)
        # 소득: detail.income_avg 또는 income_basis 명시 시 그 값, 비면 기존 use_kosis 폴백.
        want_inc = (
            (_want("income_avg", use_kosis) or _want("income_basis", use_kosis))
            if ("income_avg" in detail or "income_basis" in detail)
            else use_kosis
        )

        comp = {}
        try:
            comp = await LandInfoService().collect_comprehensive(address, pnu=pnu)
        except Exception:  # noqa: BLE001
            pass
        stats = await self._category_stats(lawd_cd)

        # ── Phase 1: 공공 인구 및 소득 데이터(SGIS, KOSIS) 연동 ──
        import asyncio

        from apps.api.app.services.market.market_models import (
            DemographicProfile,
            MacroIncomeData,
            MigrationData,
            PopulationData,
        )
        from apps.api.integrations.kosis_client import KosisClient
        from apps.api.integrations.sgis_client import SgisClient

        sgis = SgisClient()
        kosis = KosisClient()
        cur_year = str(datetime.now().year)

        # 병렬로 인구이동, 연령통계, 거시소득 호출 (옵션 선택 여부에 따라 분기)
        demographics: dict[str, Any] | None = None
        if want_pop or want_mig or want_inc:
            try:
                # use_mock=None: 클라이언트가 키 존재 여부로 실연동/폴백을 자동 결정한다.
                #   (과거 use_mock=True 하드코딩으로 키가 있어도 항상 Mock만 나오던 G1 결함 제거)
                #   키가 있으면 실데이터 시도→data_source='live', 없으면 폴백→'fallback'/'mock'/'unavailable'.
                # 통합시(수원/성남 등) 표기차 흡수: KOSIS=시 단위, SGIS=시+구.
                kosis_nm, sgis_nm = _extract_region_keys(address)

                # 항목 단위 게이팅: 미선택 항목은 provider 호출을 생략해 불필요한 외부 API
                #   호출을 줄인다(빈 메타 dict 반환 → 모델 기본값으로 unavailable 표기됨).
                async def fetch_mig():
                    if not want_mig:
                        return {"target_adm_cd": lawd_cd, "year": cur_year}
                    # I2: 인구이동은 SGIS 미제공 → KOSIS 「시군구별 이동자수」로 대상 시군구의
                    #     총전입·총전출·순이동(유입세)을 산출. 주소에서 시군구명을 추출해 식별한다.
                    od = await kosis.get_migration_od(lawd_cd[:5], cur_year, region_name=kosis_nm)
                    if od.get("data_source") == "live":
                        return od
                    # KOSIS 미확정/실패 시 SGIS 정직 폴백(가짜 금지).
                    return await sgis.get_migration_stats(lawd_cd, cur_year) if use_sgis else od
                async def fetch_pop():
                    if not want_pop:
                        return {"target_adm_cd": lawd_cd, "year": cur_year}
                    return await sgis.get_population_stats(
                        lawd_cd, cur_year, region_name=sgis_nm)
                async def fetch_inc():
                    if not want_inc:
                        return {"sigungu_cd": lawd_cd[:5], "year": cur_year}
                    return await kosis.get_macro_income_stats(
                        lawd_cd[:5], cur_year, region_name=kosis_nm)

                mig, pop, inc = await asyncio.gather(
                    fetch_mig(),
                    fetch_pop(),
                    fetch_inc(),
                    return_exceptions=True
                )
                # Pydantic 모델을 사용해 어댑터 패턴으로 데이터 표준화
                profile = DemographicProfile(
                    source_phase=1,
                    migration=MigrationData(**(mig if not isinstance(mig, Exception) else {"target_adm_cd": lawd_cd, "year": cur_year})),
                    population=PopulationData(**(pop if not isinstance(pop, Exception) else {"target_adm_cd": lawd_cd, "year": cur_year})),
                    macro_income=MacroIncomeData(**(inc if not isinstance(inc, Exception) else {"sigungu_cd": lawd_cd[:5], "year": cur_year}))
                )
                demographics = profile.model_dump()
            except Exception as e:
                logger.warning("Demographic data fetch failed", error=str(e))

        comp = comp if isinstance(comp, dict) else {}
        zone = comp.get("local_ordinance") or {}
        land_use = comp.get("land_use_plan") or {}
        basic = comp.get("land_register") or comp.get("basic") or {}
        infra = comp.get("infrastructure") or {}
        coords = comp.get("coordinates")
        # 용도지역: 여러 경로에서 견고하게 추출
        zone_type = (
            zone.get("zone_type") or land_use.get("zone_type")
            or basic.get("zone_type") or comp.get("zone_type")
        )
        official_price = None
        if comp.get("official_prices"):
            official_price = (comp["official_prices"][0] or {}).get("price_per_sqm")
        # 폴백: AutoZoningService(파이프라인 용도지역 감지기)로 보강
        if not zone_type:
            try:
                from app.services.zoning.auto_zoning_service import AutoZoningService

                az = await AutoZoningService().analyze_by_address(address)
                zone_type = az.get("zone_type")
                if not official_price and az.get("official_price_per_sqm"):
                    official_price = az.get("official_price_per_sqm")
            except Exception:  # noqa: BLE001
                pass

        # 평당가(만원/평) 요약 — 모든 시세는 면적 정규화된 평당가 기준으로 서술
        pp_by_type = {
            label: (v.get("per_pyeong") or {}).get("avg")
            for label, v in stats["trade"].items()
            if (v.get("per_pyeong") or {}).get("avg")
        }
        apt_pp = ((stats["trade"].get("아파트") or {}).get("per_pyeong") or {}).get("avg")

        # ── Phase 3: 사업 타당성 분석 (Feasibility Engine) ──
        from app.services.market.feasibility_service import FeasibilityService
        # ★면적 미상 시 기본 100평(330㎡) 임의값 제거(완성도 감사 P0 — 근거 없는 수지 날조 방지).
        #   0이면 아래 통합집계가 채우거나, 끝내 0이면 analyze가 available:False로 정직 반환.
        land_area = float(basic.get("land_area") or basic.get("area_sqm") or 0.0)

        # ── P1 다필지 통합면적 전파 ──
        # 프론트가 2필지 이상 보내면(parcels) 대표 1필지 면적이 아니라 '면적가중 통합면적'으로
        #   land_area를 덮어쓴다(예: 12필지 12,079㎡ → 대표 1,161㎡ 고착 버그 해소).
        #   ★공용 단일경유: /zoning/integrated-analysis와 동일한 ComprehensiveAnalysisService.
        #   _integrated_context(면적가중 _aggregate_integrated_zoning 재사용)를 호출 — 산식 복제 0.
        #   1필지 이하/실패면 통합 안 함(기존 단일 경로 그대로 = 무회귀).
        integrated: dict[str, Any] | None = None
        # 방어: 라우터는 parcels를 list[dict] 무스키마로 받으므로 dict 행만 통과시킨다
        #   (_integrated_context 내부 가드와 이중 안전 — 비정상 행이 섞여도 크래시 없이 무시).
        _rows = [p for p in parcels if isinstance(p, dict)] if parcels else []
        if len(_rows) >= 2:
            try:
                from app.services.land_intelligence.comprehensive_analysis_service import (
                    ComprehensiveAnalysisService,
                )

                integrated = await ComprehensiveAnalysisService()._integrated_context(_rows)
            except Exception as e:  # noqa: BLE001 — 통합집계 실패는 단일 경로로 폴백(분석 무중단)
                logger.warning("다필지 통합집계 실패 — 단일필지 경로로 폴백(graceful)", err=str(e)[:120])
                integrated = None
        if integrated and (integrated.get("total_area_sqm") or 0) > 0:
            # 통합면적으로 land_area override → 이후 feasibility·공급면적·GFA가 통합면적 기준 산출.
            land_area = float(integrated["total_area_sqm"])
            # 대표 용도지역도 통합 우세값으로 보정(미상/혼재면 기존 zone_type 유지).
            _dom = integrated.get("dominant_zone")
            if _dom and _dom not in ("mixed_review_required",):
                zone_type = _dom

        # 대표 평당가는 아파트 평당가를 우선 사용, 없으면 전체 평균 사용
        valid_pp = [v for v in pp_by_type.values() if v is not None]
        target_pp = apt_pp or (sum(valid_pp)/len(valid_pp) if valid_pp else 2000)
        # ★P0(완성도 감사): 다필지 통합집계가 이미 산출한 실효 용적률(blended_far_eff_pct,
        #   면적가중·조례반영·공용 산식)을 개략 수지에 주입 — 자체 하드코딩 표 우회 제거(SSOT).
        _blend_far = float(integrated.get("blended_far_eff_pct") or 0) if integrated else 0
        _blend_bcr = float(integrated.get("blended_bcr_eff_pct") or 0) if integrated else 0
        _far_note = (integrated.get("far_basis_note") or "") if integrated else ""
        feasibility = FeasibilityService().analyze_feasibility(
            land_area_sqm=land_area,
            zone_type=zone_type or "",
            avg_pyeong_price_manwon=target_pp,
            official_price_per_sqm=official_price or 0,
            far_pct_override=_blend_far if _blend_far > 0 else None,
            bcr_pct_override=_blend_bcr if _blend_bcr > 0 else None,
            far_basis_override=(f"다필지 통합 실효용적률(면적가중){' — ' + _far_note if _far_note else ''}"
                                if _blend_far > 0 else None),
        )

        # ── M3: 적정 분양가 산정 — 거래사례비교(1차 핵심) + 지불여력(2차 검증)·결정론 ──
        # 1차: 주변 동일종목 실거래 시세(평당가)·주변 분양가. 2차: KOSIS 소득→PIR/DSR/LTV로 수요 수용성.
        # 비교 데이터 없으면 엔진이 data_source='unavailable'로 정직 반환(가짜값 금지).
        from app.services.market.pricing_band_service import compute_fair_price
        _mi = (demographics or {}).get("macro_income") or {}
        _income_10k = _mi.get("median_income_10k") or _mi.get("avg_income_10k")
        # 실거래 평당가(만원/평) — 폴백 2000 제외, 실값만 비교가로 사용.
        _real_pp = apt_pp or (sum(valid_pp) / len(valid_pp) if valid_pp else None)
        # 대표 84㎡ 1세대 실거래 기반가(만원) = 평당가 × 25.4평(=84/3.305785)
        _trade_unit_10k = round(_real_pp * (84.0 / 3.305785)) if _real_pp else None
        # 주변 신규 분양가(청약홈) — 84㎡급(공급 100~125㎡) 분양총액 중앙값. 키/데이터 없으면 None(정직).
        _presale_10k, _presale_src = await self._nearby_presale_84_price(lawd_cd, coords)
        pricing_band = compute_fair_price(
            comparable_trade_10k=_trade_unit_10k,
            nearby_presale_10k=_presale_10k,
            annual_income_10k=_income_10k,
            trade_source="live" if _real_pp else None,
            presale_source=_presale_src,
            income_source=_mi.get("data_source"),
        )

        # ── I6: 수요기반 평형 MD 추천(가구원수 분포 → 권장 전용면적 배분)·결정론 ──
        from app.services.market.unit_mix_recommender import recommend_unit_mix
        _pop = (demographics or {}).get("population") or {}
        unit_mix_recommendation = recommend_unit_mix(
            _pop.get("household_types"),
            data_source=_pop.get("data_source"),
        )

        # ── 과제1: 실데이터 타겟 프로파일(K-Atlas 대체) — 상권(SEMAS) 신규 배선 + 5축 조립 ──
        # 상권 특성은 land_info_service 인스턴스화 패턴(CommercialAreaService)을 그대로 재사용한다.
        # 실패/미확보는 해당 축 생략(정직). 신용·카드소비는 target_profile.premium 으로 정직 강등.
        _commercial, _commercial_src = await self._commercial_area(coords)
        target_profile = _build_target_profile(demographics, _commercial, _commercial_src, infra)

        ctx = {
            "address": address,
            "zone_type": zone_type,
            "official_price_per_sqm": official_price,
            "price_basis": "평당가(만원/평) 기준으로 서술할 것",
            "apt_avg_per_pyeong_manwon": apt_pp,
            "avg_per_pyeong_by_type_manwon": pp_by_type,
            "apt_trend_per_pyeong": [
                {"ym": t["ym"], "per_pyeong_manwon": t.get("avg_per_pyeong")}
                for t in (stats.get("apt_trend") or [])
            ],
            "rent_stats_manwon": stats["rent"],
            "subway": (infra.get("nearest_subway") or {}).get("name") if isinstance(infra, dict) else None,
            "demographics": demographics,
            "feasibility": feasibility,
            # ── 과제2: 인라인 내러티브도 적정분양가·수요MD·타겟프로파일을 보게 컨텍스트 보강(재계산 0) ──
            "pricing_band": pricing_band,
            "unit_mix_recommendation": unit_mix_recommendation,
            "target_profile": target_profile,
            # AVM 참조: 별도 AVM 미산출 → 공시지가 기준 토지가액(원)을 감정 참고치로 전달(가짜 AVM 금지).
            "avm_reference": {
                "official_price_per_sqm": official_price,
                "land_appraised_total_won": (round(official_price * land_area)
                                             if (official_price and land_area) else None),
                "basis": "공시지가 기준 토지가액(참고·감정 대체 아님)",
            },
        }
        narrative = await self._narrative(ctx) if use_llm else {
            "summary": "수집된 실거래·시세 데이터 기반 시장 현황입니다. (AI 분석 미포함)",
            "opportunities": [], "risks": [], "price_trend": "", "target_persona": "AI 분석 미포함"
        }

        # ── 과제2: 시니어 정밀분석 — 전용 MarketInterpreter(감정평가/분양대행 페르소나·GROUNDING_RULE)로
        #   실거래+가격밴드+수요MD+인구+feasibility를 통합 정밀 내러티브(6키)로 산출한다.
        #   use_llm=False 또는 미가용이면 None(정직) → 렌더러가 인라인 narrative 폴백.
        senior_insight = None
        if use_llm:
            senior_insight = await self._senior_market_insight(
                address=address, zone_type=zone_type, land_area_sqm=land_area,
                stats_trade=stats["trade"], official_price=official_price,
                pricing_band=pricing_band, unit_mix=unit_mix_recommendation,
                demographics=demographics, feasibility=feasibility,
            )

        # ── 단일 소비원(raw_data + analysis) 구성 ──
        # 프론트(P3)·export(P4)가 dict 를 재가공하지 않도록 표(row 배열)로 평탄화한다.
        # 미선택 분류 처리 규칙(일관): population/income 블록은 '미선택(provider 호출 안 함)'이면
        #   키 자체를 생략하고, '선택했으나 provider 실패'면 키는 두되 data_source 로 정직 표기한다.
        raw_data: dict[str, Any] = {
            "real_estate": {
                "trade_table": _build_trade_table(stats["trade"]),
                "rent_table": _build_rent_table(stats["rent"]),
                "trend_series": _build_trend_series(stats.get("apt_trend") or []),
                "source": "국토교통부 실거래가",
                "data_source": "live",
            },
        }
        _pop_block = _build_population_block(demographics)
        if _pop_block is not None:  # 미선택이면 키 생략(정직)
            raw_data["population"] = _pop_block
        _inc_block = _build_income_block(demographics)
        if _inc_block is not None:  # 미선택이면 키 생략(정직)
            raw_data["income"] = _inc_block

        # analysis: 기존 산출을 묶기만(중복 계산 금지, 같은 객체 참조).
        analysis: dict[str, Any] = {
            "narrative": narrative,
            "feasibility": feasibility,
            "pricing_band": pricing_band,
            "unit_mix": unit_mix_recommendation,
            "target_persona": (narrative or {}).get("target_persona"),
        }

        report = {
            "generated_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
            "address": address,
            "lawd_cd": lawd_cd,
            "coordinates": coords,
            "months": stats["months"],
            "zone_type": ctx["zone_type"],
            "official_price_per_sqm": official_price,
            "trade": stats["trade"],
            "rent": stats["rent"],
            "apt_trend": stats.get("apt_trend") or [],
            "infrastructure": infra,
            "demographics": demographics,
            "narrative": narrative,
            "feasibility_analysis": feasibility,
            "pricing_band": pricing_band,
            "unit_mix_recommendation": unit_mix_recommendation,
            # ── 과제1: 실데이터 타겟 프로파일(5축·각 축 data_source 정직) ──
            "target_profile": target_profile,
            # ── 과제2: 시니어 정밀 내러티브(MarketInterpreter 6키) — 미가용/use_llm=False면 None ──
            "senior_insight": senior_insight,
            # ── 신규(하위호환): 단일 소비원 ──
            "raw_data": raw_data,
            "analysis": analysis,
        }
        # ── P1: 다필지 통합면적 메타(있으면) — 화면/PDF가 "통합 N필지(통합면적)" 표기에 사용 ──
        #   parcels 미전달/1필지면 키 자체를 생략(단일 경로 무회귀).
        if integrated and (integrated.get("total_area_sqm") or 0) > 0:
            report["integrated"] = {
                "parcel_count": integrated.get("parcel_count"),
                "total_area_sqm": integrated.get("total_area_sqm"),
                "dominant_zone": integrated.get("dominant_zone"),
                "blended_far_eff_pct": integrated.get("blended_far_eff_pct"),
                "blended_bcr_eff_pct": integrated.get("blended_bcr_eff_pct"),
                "integrated_gfa_sqm": integrated.get("integrated_gfa_sqm"),
            }
        # ── 과제2: 시니어 다전문가(감정평가·도시계획·금융) 자문 모세혈관 배선 ──
        #   실값 주입: 감정평가=공시지가 기준 토지가액 / 금융=총사업비 / 도시계획=용도·면적(프레임워크).
        self._attach_senior(
            report, feasibility,
            official_price=official_price, land_area_sqm=land_area,
            zone_type=zone_type, real_pp=_real_pp,
        )
        return report

    # ── 정적 지도 이미지(OSM 타일 합성, Pillow) ──
    @staticmethod
    def static_map_png(lat: float, lon: float, radius_m: int = 1000, zoom: int = 15,
                       w: int = 720, h: int = 440) -> bytes | None:
        """대상 좌표 중심 OSM 정적 지도 PNG(중심핀 + 반경원). 실패 시 None."""
        try:
            import math

            import httpx
            from PIL import Image, ImageDraw

            n = 2 ** zoom
            xf = (lon + 180.0) / 360.0 * n
            lat_r = math.radians(lat)
            yf = (1.0 - math.asinh(math.tan(lat_r)) / math.pi) / 2.0 * n
            cols = w // 256 + 2
            rows = h // 256 + 2
            x0 = int(xf) - cols // 2
            y0 = int(yf) - rows // 2
            canvas = Image.new("RGB", (cols * 256, rows * 256), (235, 235, 235))
            headers = {"User-Agent": "PropAI/1.0 (market report)"}
            with httpx.Client(timeout=8.0, headers=headers) as client:
                for cx in range(cols):
                    for cy in range(rows):
                        tx, ty = x0 + cx, y0 + cy
                        if tx < 0 or ty < 0 or tx >= n or ty >= n:
                            continue
                        try:
                            r = client.get(f"https://a.tile.openstreetmap.org/{zoom}/{tx}/{ty}.png")
                            if r.status_code == 200:
                                tile = Image.open(io.BytesIO(r.content)).convert("RGB")
                                canvas.paste(tile, (cx * 256, cy * 256))
                        except Exception:  # noqa: BLE001
                            continue
            # 중심 픽셀
            cpx = int((xf - x0) * 256)
            cpy = int((yf - y0) * 256)
            # 목표 크기로 중심 크롭
            left = max(0, cpx - w // 2)
            top = max(0, cpy - h // 2)
            img = canvas.crop((left, top, left + w, top + h))
            d = ImageDraw.Draw(img, "RGBA")
            ox, oy = cpx - left, cpy - top
            # 반경 원
            mpp = 156543.03392 * math.cos(lat_r) / n
            rpx = int(radius_m / mpp)
            d.ellipse([ox - rpx, oy - rpx, ox + rpx, oy + rpx], outline=(20, 184, 166, 220), width=3)
            d.ellipse([ox - rpx, oy - rpx, ox + rpx, oy + rpx], fill=(20, 184, 166, 30))
            # 중심 핀
            d.ellipse([ox - 9, oy - 9, ox + 9, oy + 9], fill=(239, 68, 68, 255), outline=(255, 255, 255, 255), width=3)
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            return buf.getvalue()
        except Exception as e:  # noqa: BLE001
            logger.warning("정적 지도 생성 실패", err=str(e)[:80])
            return None

    # ── PDF (reportlab) ──
    def to_pdf(self, rep: dict[str, Any]) -> bytes:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

        try:
            pdfmetrics.registerFont(UnicodeCIDFont("HYSMyeongJo-Medium"))
            font = "HYSMyeongJo-Medium"
        except Exception:  # noqa: BLE001
            font = "Helvetica"

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=18 * mm, bottomMargin=18 * mm)
        ss = getSampleStyleSheet()
        h1 = ParagraphStyle("h1", parent=ss["Title"], fontName=font, fontSize=20)
        h2 = ParagraphStyle("h2", parent=ss["Heading2"], fontName=font, fontSize=13, textColor=colors.HexColor("#0e7490"))
        h3 = ParagraphStyle("h3", parent=ss["Heading3"], fontName=font, fontSize=11, textColor=colors.HexColor("#0f766e"))
        body = ParagraphStyle("body", parent=ss["BodyText"], fontName=font, fontSize=10, leading=16)
        cap = ParagraphStyle("cap", parent=body, fontSize=7.5, textColor=colors.HexColor("#888888"), leading=11)
        story: list = []

        nar = rep.get("narrative") or {}
        raw = rep.get("raw_data") or {}
        tp = rep.get("target_profile") or {}
        # ── narrative 정직 가드 — 폴백/미포함 시 "AI 분석 미포함"으로 대체(무목업·정직) ──
        _ai_missing = "AI 분석 미포함(LLM 키 미설정 또는 호출 실패)"

        # ── 공용 표/내러티브 헬퍼(정의 후 사용 — 클로저이므로 상단 배치 필수) ──
        _grid_style = TableStyle([
            ("FONTNAME", (0, 0), (-1, -1), font), ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0e7490")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#cbd5e1")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f1f5f9")]),
            ("ALIGN", (1, 0), (-1, -1), "RIGHT"),
        ])

        def _simple_table(rows: list, col_widths: list):
            t = Table(rows, colWidths=col_widths)
            t.setStyle(_grid_style)
            return t

        def stat_table(title: str, data: dict):
            story.append(Paragraph(title, h3))
            rows = [["유형", "건수", "평균", "최저", "최고"]]
            for label, s in data.items():
                rows.append([label, str(s.get("count", 0)), _eok(s.get("avg", 0)), _eok(s.get("min", 0)), _eok(s.get("max", 0))])
            story.append(_simple_table(rows, [45 * mm, 25 * mm, 35 * mm, 35 * mm, 35 * mm]))
            story.append(Spacer(1, 8))

        # 매매 시세: 평당가(만원/평) 중심 + 총액 평균 병기
        def trade_table(title: str, data: dict):
            story.append(Paragraph(title, h3))
            rows = [["유형", "건수", "평당가(만원/평)", "총액 평균", "평균면적"]]
            for label, s in data.items():
                pp = (s.get("per_pyeong") or {}).get("avg", 0)
                area = s.get("avg_area_m2", 0)
                rows.append([
                    label, str(s.get("count", 0)),
                    f"{int(pp):,}만원/평" if pp else "-",
                    _eok(s.get("avg", 0)),
                    f"{area:.1f}㎡({round(area / PYEONG_SQM)}평)" if area else "-",
                ])
            story.append(_simple_table(rows, [40 * mm, 20 * mm, 42 * mm, 35 * mm, 38 * mm]))
            story.append(Spacer(1, 8))

        def _won(v):  # 만원 정수 → "○억/○만원"(0/None은 "-")
            return f"{_eok(v)}원" if v is not None and v != 0 else "-"

        # ── 데이터 출처 정직표기 헬퍼(P3-3) — mock/unavailable/None이면 가짜 라벨 대신 정직 안내 ──
        def _is_real_source(ds: Any) -> bool:
            return str(ds or "").lower() not in ("mock", "unavailable", "fallback", "")

        def _src_caption(label_source: str, ds: Any) -> str:
            if _is_real_source(ds):
                return f"출처: {label_source} ({ds})"
            return "데이터 없음(공공 API 키 미설정 — 실데이터 미연동)"

        def _narr(senior_key: str, narrative_key: str | None = None) -> None:
            """섹션 정밀 내러티브 1~2문장 주입(senior_insight 우선·인라인 폴백). 표만 있던 섹션에 해석 부여."""
            txt = _insight_text(rep, senior_key, narrative_key)
            if txt:
                story.append(Paragraph(_esc(txt), body))
                story.append(Spacer(1, 4))

        # ══ 표지 ══
        story.append(Paragraph("시장조사보고서", h1))
        # address·generated_at 는 사용자/엔진 동적 입력이라 esc(주소에 '<','&' 가능 → Paragraph 크래시 차단).
        story.append(Paragraph(
            f"{_esc(rep['address'])} · 생성 {_esc(rep['generated_at'])} · 최근 {len(rep['months'])}개월", body))
        # ── 다필지 통합면적 배너(있으면) — "통합 N필지" 표기로 대표 1필지 오인 방지 ──
        _intg = rep.get("integrated") or {}
        if _intg.get("parcel_count") and (_intg.get("total_area_sqm") or 0) > 0:
            _ta = _intg["total_area_sqm"]
            story.append(Paragraph(
                f"통합 {_intg['parcel_count']}필지 · 통합 대지면적 "
                f"{_ta:,.0f}㎡({round(_ta / PYEONG_SQM):,}평) 기준 분석", body))
        story.append(Spacer(1, 10))

        # ══ Executive Summary(두괄식) — 등급 + 핵심 KPI + 투자의견 ══
        es = _exec_summary(rep)
        story.append(Paragraph("핵심 요약 (Executive Summary)", h2))
        story.append(Paragraph(
            f"시장 강도: <b>{_esc(es['market_strength']['grade'])}</b>"
            f"({_esc(es['market_strength']['label'])}) · "
            f"사업성 등급: <b>{_esc(es['business_grade']['grade'])}</b>"
            f"({_esc(es['business_grade']['label'])})", body))
        _kpi = es["kpi"]
        _apt, _nm, _fp2, _roi2 = (_kpi["apt_per_pyeong_manwon"], _kpi["net_migration"],
                                  _kpi["fair_price_10k"], _kpi["roi_percent"])
        krows = [["핵심 지표", "값"],
                 ["아파트 평당시세", f"{int(_apt):,}만원/평" if _apt else "-"],
                 ["순유입(전입-전출)", f"{_nm:,}명" if _nm is not None else "-"],
                 ["적정 분양가(84㎡)", f"{_eok(_fp2)}원" if _fp2 else "-"],
                 ["개략 ROI", f"{_roi2:.1f}%" if _roi2 is not None else "-"]]
        story.append(_simple_table(krows, [70 * mm, 100 * mm]))
        # opinion·expert_opinion 은 엔진/LLM 동적 문자열이라 esc.
        story.append(Paragraph(f"투자의견: {_esc(es['opinion'])}", body))
        if es.get("expert_opinion"):
            story.append(Paragraph(_esc(es["expert_opinion"]), cap))
        story.append(Spacer(1, 10))

        # ══ 1. 시장 개요 ══
        story.append(Paragraph("1. 시장 개요", h2))
        # senior_insight(market_overview) 우선 → 인라인 narrative.summary → AI미포함(정직).
        _ov = _insight_text(rep, "market_overview", "summary") or _ai_missing
        story.append(Paragraph(_esc(_ov), body))
        if rep.get("zone_type") or rep.get("official_price_per_sqm"):
            opp = rep.get("official_price_per_sqm")
            # ★공시지가 0/None 가드: 0(정수)은 falsy이나 명시 가드로 "0만" 오표시 차단 → "-".
            opp_txt = _eok(opp / 10000) if (opp and opp > 0) else "-"
            story.append(Paragraph(
                f"용도지역: {_esc(rep.get('zone_type') or '-')} · 공시지가(㎡): {opp_txt}", body))
        story.append(Spacer(1, 8))

        # ══ 2. 수요 분석(인구 + 소득) ══
        story.append(Paragraph("2. 수요 분석 (인구·소득)", h2))
        pop = raw.get("population")
        if pop and not _is_real_source(pop.get("data_source")):
            story.append(Paragraph("데이터 없음(공공 API 키 미설정 — 인구통계 실데이터 미연동)", cap))
        elif pop:
            story.append(Paragraph("인구 규모·분포", h3))
            summ = pop.get("summary") or {}
            srows = [["구분", "값"],
                     ["총인구", f"{summ.get('total_population'):,}명" if summ.get("total_population") else "데이터 없음"],
                     ["가구수", f"{summ.get('household_count'):,}가구" if summ.get("household_count") else "데이터 없음"],
                     ["평균 가구원수", f"{summ.get('avg_household_size')}명" if summ.get("avg_household_size") else "데이터 없음"]]
            story.append(_simple_table(srows, [70 * mm, 100 * mm]))
            story.append(Spacer(1, 4))
            age = pop.get("age_distribution") or []
            if age:
                arows = [["연령대", "인구수"]] + [[a.get("label", "-"), f"{a.get('count'):,}" if a.get("count") else "-"] for a in age]
                story.append(Paragraph("연령 분포", body))
                story.append(_simple_table(arows, [85 * mm, 85 * mm]))
                story.append(Spacer(1, 4))
            ht = pop.get("household_types") or []
            if ht:
                hrows = [["가구원수", "비율(%)"]] + [[h.get("label", "-"), f"{h.get('ratio')}%(추정)" if h.get("ratio") is not None else "-"] for h in ht]
                story.append(Paragraph("가구원수 분포", body))
                story.append(_simple_table(hrows, [85 * mm, 85 * mm]))
                story.append(Spacer(1, 4))
            mig = pop.get("migration") or {}
            mrows = [["순유입(전입-전출)", "전입", "전출"],
                     [f"{mig.get('net_migration'):,}" if mig.get("net_migration") is not None else "데이터 없음",
                      f"{mig.get('total_inflow'):,}" if mig.get("total_inflow") is not None else "-",
                      f"{mig.get('total_outflow'):,}" if mig.get("total_outflow") is not None else "-"]]
            story.append(Paragraph("인구 이동", body))
            story.append(_simple_table(mrows, [60 * mm, 55 * mm, 55 * mm]))
            story.append(Paragraph(_esc(_src_caption(pop.get("source", "-"), pop.get("data_source"))), cap))
            story.append(Spacer(1, 6))

        inc = raw.get("income")
        if inc and not _is_real_source(inc.get("data_source")):
            story.append(Paragraph("소득 수준: 데이터 없음(공공 API 키 미설정 — 소득 실데이터 미연동)", cap))
        elif inc:
            story.append(Paragraph("소득 수준", h3))
            avg = inc.get("avg_income_10k")
            med = inc.get("median_income_10k")
            irows = [["구분", "연소득"],
                     ["평균 소득", f"{_eok(avg)}원" if avg else "데이터 없음"],
                     ["중위 소득" + ("(추정)" if inc.get("median_estimated") else ""), f"{_eok(med)}원" if med else "데이터 없음"]]
            story.append(_simple_table(irows, [80 * mm, 90 * mm]))
            story.append(Paragraph(_esc(_src_caption(inc.get("source", "-"), inc.get("data_source"))), cap))
        # 타겟 프로파일 요약(주력 연령·가구·소득분위) — 실데이터 축만 정직 표기.
        _tp_bits: list[str] = []
        _pa = tp.get("primary_age") or {}
        _ph = tp.get("primary_household") or {}
        _it = tp.get("income_tier") or {}
        if _pa.get("band"):
            _tp_bits.append(f"주력 연령대 {_pa['band']}")
        if _ph.get("type"):
            _tp_bits.append(f"주력 가구 {_ph['type']}(추정)")
        if _it.get("tier_label"):
            _tp_bits.append(f"소득 수준 {_it['tier_label']}")
        if _tp_bits:
            story.append(Paragraph(_esc("타겟 프로파일: " + " · ".join(_tp_bits)), body))
        _narr("comparable_analysis")
        story.append(Spacer(1, 8))

        # ══ 3. 가격·실거래(매매·추이·전월세·적정분양가) ══
        story.append(Paragraph("3. 가격·실거래", h2))
        trade_table("매매 시세 (유형별 · 평당가 기준)", rep.get("trade") or {})
        # 시세 추이 차트(아파트 월별 평당가)
        trend = [t for t in (rep.get("apt_trend") or []) if t.get("avg_per_pyeong") or t.get("avg")]
        if trend:
            from reportlab.graphics.charts.barcharts import VerticalBarChart
            from reportlab.graphics.shapes import Drawing

            story.append(Paragraph("매매 시세 추이 (아파트 월별 평당가, 만원/평)", h3))
            d = Drawing(440, 170)
            bc = VerticalBarChart()
            bc.x = 40; bc.y = 25; bc.width = 360; bc.height = 120
            bc.data = [[int(t.get("avg_per_pyeong") or t.get("avg") or 0) for t in trend]]
            bc.categoryAxis.categoryNames = [f"{int(t['ym'][4:6])}월" for t in trend]
            bc.categoryAxis.labels.fontName = font
            bc.valueAxis.labels.fontName = font
            bc.barWidth = 14
            bc.bars[0].fillColor = colors.HexColor("#0e7490")
            bc.valueAxis.valueMin = 0
            d.add(bc)
            story.append(d)
            story.append(Spacer(1, 8))
        stat_table("전월세 보증금 (유형별)", rep.get("rent") or {})
        # 적정 분양가(★사업타당성 앞으로 — 매출/가격이 사업성 산정의 입력) — 있을 때만.
        pb = (rep.get("pricing_band") or {})
        if pb and pb.get("data_source") not in (None, "unavailable") and pb.get("fair_price_10k"):
            story.append(Paragraph("적정 분양가 (거래사례비교)", h3))
            story.append(Paragraph(
                f"적정 분양가: {_eok(pb['fair_price_10k'])}원 · "
                f"지불여력 판정: {_esc(pb.get('affordability_verdict', '-'))}", body))
            story.append(Paragraph(_esc(pb.get("note", "")), cap))
            story.append(Spacer(1, 4))
        # 가격 동향 정밀 내러티브(senior price_trend_analysis 우선 → 인라인 price_trend).
        _narr("price_trend_analysis", "price_trend")
        story.append(Spacer(1, 6))

        # ══ 4. 입지 분석(지도 + 상권 + 인프라) ══
        story.append(Paragraph("4. 입지 분석", h2))
        coords = rep.get("coordinates") or {}
        if coords.get("lat") and coords.get("lon"):
            png = self.static_map_png(coords["lat"], coords["lon"], 1000)
            if png:
                from reportlab.platypus import Image as RLImage

                story.append(Paragraph("대상지 위치 (반경 1km)", h3))
                story.append(RLImage(io.BytesIO(png), width=165 * mm, height=101 * mm))
                story.append(Spacer(1, 6))
        # 상권 특성(SEMAS) — 실데이터만, 미확보는 정직 표기.
        comm = tp.get("commercial") or {}
        if comm.get("data_source") == "live":
            _cat = ", ".join(f"{c.get('category')}({c.get('count')})"
                             for c in (comm.get("category_distribution") or [])[:4])
            story.append(Paragraph(
                f"상권 활성도: {_esc(str(comm.get('grade') or '-'))}등급"
                f"(점수 {comm.get('vitality_score')}) · 점포 {comm.get('total_stores'):,}개"
                + (f" · 주요업종 {_esc(_cat)}" if _cat else ""), body))
        else:
            story.append(Paragraph("상권(SEMAS) 데이터 미확보 — 키 미설정/좌표 부재(정직 표기)", cap))
        loc = tp.get("location") or {}
        if loc.get("nearest_subway"):
            _dist = f" ({loc['subway_distance_m']}m)" if loc.get("subway_distance_m") else ""
            story.append(Paragraph(f"최인접 지하철: {_esc(str(loc['nearest_subway']))}{_dist}"
                                   + (f" · 학교 {loc['school_count']}개소" if loc.get("school_count") else ""), body))
        story.append(Spacer(1, 8))

        # ══ 5. 사업 타당성(Feasibility · 개략 추정) ══
        story.append(Paragraph("5. 사업 타당성 (Feasibility · 개략 추정)", h2))
        fe = rep.get("feasibility_analysis") or {}
        _fin = fe.get("financials") or {}
        _mass = fe.get("massing") or {}
        if _fin and not fe.get("error"):
            _roi = _fin.get("roi_percent")
            _la = _mass.get("land_area_sqm")
            _gfa = _mass.get("gfa_sqm")
            _land_txt = (f"{_la:,.0f}㎡({round(_la / PYEONG_SQM):,}평)" if _la else "-")
            _gfa_txt = (f"{_gfa:,.0f}㎡({round(_mass.get('gfa_pyeong') or 0):,}평)" if _gfa else "-")
            frows = [
                ["구분", "값"],
                ["ROI(투자수익률)", f"{_roi:.1f}%" if _roi is not None else "-"],
                ["총사업비", _won(_fin.get("total_cost_10k"))],
                ["예상 분양수익", _won(_fin.get("total_revenue_10k"))],
                ["세전 순수익", _won(_fin.get("net_profit_10k"))],
                ["NPV(순현재가치)", _won(_fin.get("npv_10k"))],
                ["대지면적", _land_txt],
                ["건축가능 연면적", _gfa_txt],
            ]
            story.append(_simple_table(frows, [70 * mm, 100 * mm]))
            story.append(Paragraph(_esc((fe.get("assumptions") or {}).get("note", "")), cap))
            story.append(Spacer(1, 6))
            # 평형 MD(수요기반 전용면적 배분) — 데이터 있으면 표로(unavailable이면 정직 표기).
            um = rep.get("unit_mix_recommendation") or {}
            mix = um.get("recommended_mix") or {}
            if mix and um.get("data_source") not in (None, "unavailable"):
                story.append(Paragraph("권장 평형 배분(수요기반 MD)", body))
                mrows = [["평형대", "권장 배분(%)"]] + [[k, f"{v}%"] for k, v in mix.items()]
                story.append(_simple_table(mrows, [85 * mm, 85 * mm]))
                story.append(Paragraph(_esc(um.get("rationale", "")), cap))
            elif um.get("data_source") == "unavailable":
                story.append(Paragraph(
                    "권장 평형 배분: 데이터 없음(인구/가구 분석 미선택 — 가구원수 분포 필요)", cap))
            # 사업성 정밀 해석(senior investment_insight) — 표 뒤 수치→해석.
            _narr("investment_insight")
        else:
            story.append(Paragraph(
                "개략 수지 산출 불가(면적/용도지역 근거 부족) — 무목업 정직 표기", cap))
        story.append(Spacer(1, 8))

        # ══ 6. 리스크 요인 ══
        story.append(Paragraph("6. 리스크 요인", h2))
        _rf = _insight_text(rep, "risk_factors")
        _risks = [r for r in (nar.get("risks") or []) if str(r).strip()]
        if _rf:
            story.append(Paragraph(_esc(_rf), body))
        if _risks:
            for r in _risks:
                story.append(Paragraph(f"· {_esc(r)}", body))
        if not _rf and not _risks:
            story.append(Paragraph(_ai_missing, cap))
        story.append(Spacer(1, 8))

        # ══ 7. 결론 및 권고 ══
        story.append(Paragraph("7. 결론 및 권고", h2))
        story.append(Paragraph(f"투자의견: {_esc(es['opinion'])}", body))
        _narr("timing_recommendation")
        _narr("investment_insight")
        _opps = [o for o in (nar.get("opportunities") or []) if str(o).strip()]
        if _opps:
            story.append(Paragraph("기회 요인", h3))
            for o in _opps:
                story.append(Paragraph(f"· {_esc(o)}", body))
        story.append(Spacer(1, 8))

        # ══ 8. 부록 · 출처 및 면책 ══
        story.append(Paragraph("8. 부록 · 출처 및 면책", h2))
        re_block = raw.get("real_estate") or {}
        story.append(Paragraph(_esc(f"실거래 출처: {re_block.get('source', '국토교통부 실거래가')}"), cap))
        prem = tp.get("premium") or {}
        if prem.get("note"):
            story.append(Paragraph(_esc(prem["note"]), cap))
        story.append(Spacer(1, 6))
        disc = ParagraphStyle("disc", parent=body, fontSize=7.5, textColor=colors.HexColor("#888888"), leading=11)
        story.append(Paragraph(DISCLAIMER_TEXT, disc))

        doc.build(story)
        return buf.getvalue()

    # ── PPTX (python-pptx) ──
    def to_pptx(self, rep: dict[str, Any]) -> bytes:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        accent = RGBColor(0x0E, 0x74, 0x90)
        ink = RGBColor(0x0F, 0x17, 0x2A)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)

        def _fill(shape, rgb):
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb
            shape.line.fill.background()

        def brand_footer(s):
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), Inches(13.33), Inches(0.4))
            _fill(bar, accent)
            tf = bar.text_frame
            tf.margin_top = Pt(2)
            tf.text = "사통팔땅 · AI 부동산 인텔리전스   |   시장조사보고서"
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.color.rgb = WHITE

        def header_bar(s, title: str):
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.1))
            _fill(bar, accent)
            tf = bar.text_frame
            tf.margin_left = Inches(0.6)
            tf.word_wrap = True
            tf.text = title
            tf.paragraphs[0].font.size = Pt(26)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = WHITE

        def title_slide():
            s = prs.slides.add_slide(prs.slide_layouts[6])
            # 브랜드 풀블리드 배경
            bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
            _fill(bg, ink)
            band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.0), Inches(13.33), Inches(0.12))
            _fill(band, accent)
            brand = s.shapes.add_textbox(Inches(0.9), Inches(0.7), Inches(11), Inches(0.6)).text_frame
            brand.text = "사통팔땅  ·  AI 부동산 인텔리전스"
            brand.paragraphs[0].font.size = Pt(16); brand.paragraphs[0].font.color.rgb = accent; brand.paragraphs[0].font.bold = True
            tb = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.7), Inches(2.5)).text_frame
            tb.text = "시장조사보고서"
            tb.paragraphs[0].font.size = Pt(52); tb.paragraphs[0].font.bold = True; tb.paragraphs[0].font.color.rgb = WHITE
            p = tb.add_paragraph()
            p.text = f"{rep['address']}"
            p.font.size = Pt(22); p.font.color.rgb = WHITE
            p2 = tb.add_paragraph()
            p2.text = f"생성 {rep['generated_at']} · 최근 {len(rep['months'])}개월 · 실거래 기반"
            p2.font.size = Pt(14); p2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
            # 면책 고지
            disc = s.shapes.add_textbox(Inches(0.9), Inches(6.7), Inches(11.7), Inches(0.6)).text_frame
            disc.word_wrap = True
            disc.text = DISCLAIMER_TEXT
            disc.paragraphs[0].font.size = Pt(9); disc.paragraphs[0].font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

        def map_slide():
            coords = rep.get("coordinates") or {}
            if not (coords.get("lat") and coords.get("lon")):
                return
            png = self.static_map_png(coords["lat"], coords["lon"], 1000)
            if not png:
                return
            s = prs.slides.add_slide(prs.slide_layouts[6])
            header_bar(s, "대상지 위치 (반경 1km)")
            s.shapes.add_picture(io.BytesIO(png), Inches(2.6), Inches(1.4), height=Inches(5.4))
            brand_footer(s)

        def text_slide(title: str, lines: list[str]):
            s = prs.slides.add_slide(prs.slide_layouts[6])
            header_bar(s, title)
            bodytf = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2)).text_frame
            bodytf.word_wrap = True
            for i, ln in enumerate(lines or ["-"]):
                para = bodytf.paragraphs[0] if i == 0 else bodytf.add_paragraph()
                para.text = ln
                para.font.size = Pt(16); para.font.color.rgb = ink
                para.space_after = Pt(8)
            brand_footer(s)

        def table_slide(title: str, data: dict, pp: bool = False):
            s = prs.slides.add_slide(prs.slide_layouts[6])
            header_bar(s, title)
            rows = len(data) + 1
            ncol = 4 if pp else 5
            tbl = s.shapes.add_table(rows, ncol, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5 * rows)).table
            hdr = ["유형", "건수", "평당가(만원/평)", "총액 평균"] if pp else ["유형", "건수", "평균", "최저", "최고"]
            for c, h in enumerate(hdr):
                cell = tbl.cell(0, c)
                cell.text = h
                cell.fill.solid(); cell.fill.fore_color.rgb = accent
                cell.text_frame.paragraphs[0].font.color.rgb = WHITE
                cell.text_frame.paragraphs[0].font.bold = True
            for r, (label, st) in enumerate(data.items(), start=1):
                tbl.cell(r, 0).text = label
                tbl.cell(r, 1).text = str(st.get("count", 0))
                if pp:
                    ppv = (st.get("per_pyeong") or {}).get("avg", 0)
                    tbl.cell(r, 2).text = f"{int(ppv):,}만원/평" if ppv else "-"
                    tbl.cell(r, 3).text = _eok(st.get("avg", 0))
                else:
                    tbl.cell(r, 2).text = _eok(st.get("avg", 0))
                    tbl.cell(r, 3).text = _eok(st.get("min", 0))
                    tbl.cell(r, 4).text = _eok(st.get("max", 0))
            brand_footer(s)

        def chart_slide(title: str, trend: list[dict[str, Any]]):
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE

            s = prs.slides.add_slide(prs.slide_layouts[6])
            header_bar(s, title)
            cd = CategoryChartData()
            cd.categories = [f"{int(x['ym'][4:6])}월" for x in trend]
            cd.add_series("아파트 평당가(만원/평)", [int(x.get("avg_per_pyeong") or x.get("avg") or 0) for x in trend])
            s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5), cd)
            brand_footer(s)

        def kv_table_slide(title: str, header: list[str], data_rows: list[list[str]], caption: str | None = None):
            """2~N열 일반 표 슬라이드(인구·소득 등). data_rows 는 문자열 행 배열."""
            s = prs.slides.add_slide(prs.slide_layouts[6])
            header_bar(s, title)
            rows = len(data_rows) + 1
            ncol = len(header)
            tbl = s.shapes.add_table(rows, ncol, Inches(0.8), Inches(1.5), Inches(11.7), Inches(min(0.5 * rows, 5))).table
            for c, h in enumerate(header):
                cell = tbl.cell(0, c)
                cell.text = h
                cell.fill.solid(); cell.fill.fore_color.rgb = accent
                cell.text_frame.paragraphs[0].font.color.rgb = WHITE
                cell.text_frame.paragraphs[0].font.bold = True
            for r, row in enumerate(data_rows, start=1):
                for c, val in enumerate(row):
                    tbl.cell(r, c).text = str(val)
            if caption:
                capbox = s.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5)).text_frame
                capbox.word_wrap = True
                capbox.text = caption
                capbox.paragraphs[0].font.size = Pt(9)
                capbox.paragraphs[0].font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
            brand_footer(s)

        nar = rep.get("narrative") or {}
        trend = [x for x in (rep.get("apt_trend") or []) if x.get("avg_per_pyeong") or x.get("avg")]
        raw = rep.get("raw_data") or {}
        tp = rep.get("target_profile") or {}
        _ai_missing = "AI 분석 미포함(LLM 키 미설정 또는 호출 실패)"
        es = _exec_summary(rep)

        # ══ 표지 ══
        title_slide()

        # ══ Executive Summary(두괄식) — 등급 + 핵심 KPI + 투자의견 ══
        _kpi = es["kpi"]
        _es_lines = [
            f"시장 강도: {es['market_strength']['grade']}({es['market_strength']['label']})  ·  "
            f"사업성 등급: {es['business_grade']['grade']}({es['business_grade']['label']})",
            (f"아파트 평당시세: {int(_kpi['apt_per_pyeong_manwon']):,}만원/평"
             if _kpi["apt_per_pyeong_manwon"] else "아파트 평당시세: -"),
            (f"순유입(전입-전출): {_kpi['net_migration']:,}명"
             if _kpi["net_migration"] is not None else "순유입(전입-전출): -"),
            (f"적정 분양가(84㎡): {_eok(_kpi['fair_price_10k'])}원"
             if _kpi["fair_price_10k"] else "적정 분양가(84㎡): -"),
            (f"개략 ROI: {_kpi['roi_percent']:.1f}%"
             if _kpi["roi_percent"] is not None else "개략 ROI: -"),
            f"투자의견: {es['opinion']}",
        ]
        if es.get("expert_opinion"):
            _es_lines.append(es["expert_opinion"])
        text_slide("핵심 요약 (Executive Summary)", _es_lines)

        # ══ 1. 시장 개요 ══
        _ov = _insight_text(rep, "market_overview", "summary") or _ai_missing
        text_slide("1. 시장 개요", [_ov, f"용도지역: {rep.get('zone_type') or '-'}"])

        # ══ 2. 수요 분석(인구 + 소득 + 타겟 프로파일) ══
        pop = raw.get("population")
        if pop:
            summ = pop.get("summary") or {}
            prows = [
                ["총인구", f"{summ.get('total_population'):,}명" if summ.get("total_population") else "데이터 없음"],
                ["가구수", f"{summ.get('household_count'):,}가구" if summ.get("household_count") else "데이터 없음"],
                ["평균 가구원수", f"{summ.get('avg_household_size')}명" if summ.get("avg_household_size") else "데이터 없음"],
            ]
            for a in (pop.get("age_distribution") or []):
                prows.append([f"연령 {a.get('label', '-')}", f"{a.get('count'):,}명" if a.get("count") else "-"])
            for h in (pop.get("household_types") or []):
                prows.append([f"가구원 {h.get('label', '-')}", f"{h.get('ratio')}%(추정)" if h.get("ratio") is not None else "-"])
            mig = pop.get("migration") or {}
            prows.append(["순유입(전입-전출)", f"{mig.get('net_migration'):,}명" if mig.get("net_migration") is not None else "데이터 없음"])
            kv_table_slide("2. 수요 분석 (인구)", ["구분", "값"], prows,
                           caption=f"출처: {pop.get('source', '-')} ({pop.get('data_source', '-')})")
        inc = raw.get("income")
        if inc:
            avg = inc.get("avg_income_10k")
            med = inc.get("median_income_10k")
            irows = [
                ["평균 소득", f"{_eok(avg)}원" if avg else "데이터 없음"],
                ["중위 소득" + ("(추정)" if inc.get("median_estimated") else ""), f"{_eok(med)}원" if med else "데이터 없음"],
            ]
            kv_table_slide("2. 수요 분석 (소득)", ["구분", "연소득"], irows,
                           caption=f"출처: {inc.get('source', '-')} ({inc.get('data_source', '-')})")
        # 타겟 프로파일 + 수요 정밀 해석(항상 1슬라이드 — 섹션 2 존재 보장·정직)
        _tp_bits: list[str] = []
        _pa = tp.get("primary_age") or {}
        _ph = tp.get("primary_household") or {}
        _it = tp.get("income_tier") or {}
        if _pa.get("band"):
            _tp_bits.append(f"주력 연령대: {_pa['band']}")
        if _ph.get("type"):
            _tp_bits.append(f"주력 가구: {_ph['type']}(추정)")
        if _it.get("tier_label"):
            _tp_bits.append(f"소득 수준: {_it['tier_label']}")
        _ca = _insight_text(rep, "comparable_analysis")
        if _ca:
            _tp_bits.append(_ca)
        if not _tp_bits:
            _tp_bits = ["타겟 프로파일 데이터 미확보(인구/소득 분석 미선택 — 정직 표기)"]
        text_slide("2. 수요 분석 (타겟 프로파일)", _tp_bits)

        # ══ 3. 가격·실거래(매매·추이·전월세·적정분양가) ══
        table_slide("3. 가격·실거래 — 매매 시세 (유형별 · 평당가)", rep.get("trade") or {}, pp=True)
        if trend:
            chart_slide("3. 가격·실거래 — 매매 시세 추이 (아파트 월별 평당가)", trend)
        table_slide("3. 가격·실거래 — 전월세 보증금 (유형별)", rep.get("rent") or {})
        # 적정 분양가(★사업타당성 앞으로 — PPTX 누락 해소) — 있을 때만.
        pb = rep.get("pricing_band") or {}
        if pb.get("data_source") not in (None, "unavailable") and pb.get("fair_price_10k"):
            text_slide("3. 가격·실거래 — 적정 분양가 (거래사례비교)", [
                f"적정 분양가: {_eok(pb['fair_price_10k'])}원",
                f"지불여력 판정: {pb.get('affordability_verdict', '-')}",
                pb.get("note", ""),
            ])
        _pt = _insight_text(rep, "price_trend_analysis", "price_trend")
        if _pt:
            text_slide("3. 가격·실거래 — 가격 동향 분석", [_pt])

        # ══ 4. 입지 분석(지도 + 상권 + 인프라) ══
        map_slide()
        comm = tp.get("commercial") or {}
        loc = tp.get("location") or {}
        _loc_lines: list[str] = []
        if comm.get("data_source") == "live":
            _cat = ", ".join(f"{c.get('category')}({c.get('count')})"
                             for c in (comm.get("category_distribution") or [])[:4])
            _loc_lines.append(
                f"상권 활성도: {comm.get('grade')}등급(점수 {comm.get('vitality_score')}) · "
                f"점포 {comm.get('total_stores'):,}개" + (f" · 주요업종 {_cat}" if _cat else ""))
        else:
            _loc_lines.append("상권(SEMAS) 데이터 미확보 — 키 미설정/좌표 부재(정직 표기)")
        if loc.get("nearest_subway"):
            _dist = f" ({loc['subway_distance_m']}m)" if loc.get("subway_distance_m") else ""
            _loc_lines.append(f"최인접 지하철: {loc['nearest_subway']}{_dist}"
                              + (f" · 학교 {loc['school_count']}개소" if loc.get("school_count") else ""))
        text_slide("4. 입지 분석", _loc_lines)

        # ══ 5. 사업 타당성(Feasibility) — PPTX 누락 해소 ══
        fe = rep.get("feasibility_analysis") or {}
        _fin = fe.get("financials") or {}
        if _fin and not fe.get("error"):
            _roi = _fin.get("roi_percent")

            def _won(v):
                return f"{_eok(v)}원" if v not in (None, 0) else "-"

            frows = [
                ["ROI(투자수익률)", f"{_roi:.1f}%" if _roi is not None else "-"],
                ["총사업비", _won(_fin.get("total_cost_10k"))],
                ["예상 분양수익", _won(_fin.get("total_revenue_10k"))],
                ["세전 순수익", _won(_fin.get("net_profit_10k"))],
                ["NPV(순현재가치)", _won(_fin.get("npv_10k"))],
            ]
            kv_table_slide("5. 사업 타당성 (Feasibility · 개략 추정)", ["구분", "값"], frows,
                           caption=(fe.get("assumptions") or {}).get("note", ""))
            _ii = _insight_text(rep, "investment_insight")
            if _ii:
                text_slide("5. 사업 타당성 — 정밀 해석", [_ii])
        else:
            text_slide("5. 사업 타당성 (Feasibility)",
                       ["개략 수지 산출 불가(면적/용도지역 근거 부족) — 무목업 정직 표기"])

        # ══ 6. 리스크 요인 ══
        _rf = _insight_text(rep, "risk_factors")
        _risk_lines = ([_rf] if _rf else []) + [f"· {r}" for r in (nar.get("risks") or []) if str(r).strip()]
        if not _risk_lines:
            _risk_lines = [_ai_missing]
        text_slide("6. 리스크 요인", _risk_lines)

        # ══ 7. 결론 및 권고 ══
        _concl = [f"투자의견: {es['opinion']}"]
        for _k in ("timing_recommendation", "investment_insight"):
            _v = _insight_text(rep, _k)
            if _v:
                _concl.append(_v)
        _opps = [f"· {o}" for o in (nar.get("opportunities") or []) if str(o).strip()]
        if _opps:
            _concl.append("[기회 요인]")
            _concl.extend(_opps)
        text_slide("7. 결론 및 권고", _concl)

        # ══ 8. 부록 · 출처 및 면책 ══
        re_block = raw.get("real_estate") or {}
        _app = [f"실거래 출처: {re_block.get('source', '국토교통부 실거래가')}"]
        prem = tp.get("premium") or {}
        if prem.get("note"):
            _app.append(prem["note"])
        _app.append(DISCLAIMER_TEXT)
        text_slide("8. 부록 · 출처 및 면책", _app)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ── DOCX (python-docx) ──
    def to_docx(self, rep: dict[str, Any]) -> bytes:
        """시장조사보고서 Word(.docx) 생성. PDF/PPTX 와 동일 데이터(raw_data·analysis)·동일 섹션 구성.

        가짜 데이터 금지: 값이 없으면 '-'/'데이터 없음', 미선택 분류(키 없음)는 섹션 자체를 정직 표기.
        """
        from docx import Document
        from docx.shared import Pt, RGBColor

        doc = Document()

        # months 는 실제 rep 에선 리스트(개월 코드 배열), 스모크에선 정수일 수 있어 둘 다 허용.
        months = rep.get("months")
        months_n = len(months) if isinstance(months, (list, tuple)) else (months or 0)

        ACCENT = RGBColor(0x0E, 0x74, 0x90)

        def _h(text: str, level: int = 1):
            h = doc.add_heading(text, level=level)
            for run in h.runs:
                run.font.color.rgb = ACCENT
            return h

        def _table(header: list[str], data_rows: list[list[str]]):
            """기본 스타일 표. 첫 행=헤더. 셀 값은 문자열로 강제(가짜 금지·None→'-')."""
            t = doc.add_table(rows=1, cols=len(header))
            try:
                t.style = "Light Grid Accent 1"
            except Exception:  # noqa: BLE001  # 스타일 없으면 기본 스타일 유지
                pass
            for c, htxt in enumerate(header):
                t.rows[0].cells[c].text = htxt
            for row in data_rows:
                cells = t.add_row().cells
                for c, val in enumerate(row):
                    cells[c].text = "-" if val is None else str(val)
            doc.add_paragraph()
            return t

        def _caption(text: str):
            p = doc.add_paragraph(text)
            run = p.runs[0] if p.runs else p.add_run(text)
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

        # ── ① 표지 ──
        title = doc.add_heading("시장조사보고서", level=0)
        for run in title.runs:
            run.font.color.rgb = ACCENT
        doc.add_paragraph(f"대상지: {rep.get('address') or '-'}")
        doc.add_paragraph(f"생성일: {rep.get('generated_at') or '-'}")
        doc.add_paragraph(f"분석기간: 최근 {months_n}개월")
        doc.add_paragraph()

        nar = rep.get("narrative") or {}
        raw = rep.get("raw_data") or {}
        re_block = raw.get("real_estate") or {}

        tp = rep.get("target_profile") or {}
        _ai_missing = "AI 분석 미포함(LLM 키 미설정 또는 호출 실패)"
        es = _exec_summary(rep)

        # ── Executive Summary(두괄식) — 등급 + 핵심 KPI + 투자의견 ──
        _h("핵심 요약 (Executive Summary)", 1)
        doc.add_paragraph(
            f"시장 강도: {es['market_strength']['grade']}({es['market_strength']['label']}) · "
            f"사업성 등급: {es['business_grade']['grade']}({es['business_grade']['label']})")
        _kpi = es["kpi"]
        _table(["핵심 지표", "값"], [
            ["아파트 평당시세", f"{int(_kpi['apt_per_pyeong_manwon']):,}만원/평" if _kpi["apt_per_pyeong_manwon"] else "-"],
            ["순유입(전입-전출)", f"{_kpi['net_migration']:,}명" if _kpi["net_migration"] is not None else "-"],
            ["적정 분양가(84㎡)", f"{_eok(_kpi['fair_price_10k'])}원" if _kpi["fair_price_10k"] else "-"],
            ["개략 ROI", f"{_kpi['roi_percent']:.1f}%" if _kpi["roi_percent"] is not None else "-"],
        ])
        doc.add_paragraph(f"투자의견: {es['opinion']}")
        if es.get("expert_opinion"):
            _caption(es["expert_opinion"])

        # ── 1. 시장 개요 ──
        _h("1. 시장 개요", 1)
        doc.add_paragraph(_insight_text(rep, "market_overview", "summary") or _ai_missing)
        if rep.get("zone_type") or rep.get("official_price_per_sqm"):
            opp = rep.get("official_price_per_sqm")
            doc.add_paragraph(
                f"용도지역: {rep.get('zone_type') or '-'} · "
                f"공시지가(㎡): {_eok((opp or 0) / 10000) + '원' if opp else '-'}"
            )

        # ── 2. 수요 분석 (인구 + 소득 + 타겟 프로파일) ──
        _h("2. 수요 분석 (인구·소득)", 1)
        pop = raw.get("population")
        if pop:
            _h("인구 규모·분포", 2)
            summ = pop.get("summary") or {}
            _table(["구분", "값"], [
                ["총인구", f"{summ.get('total_population'):,}명" if summ.get("total_population") else "데이터 없음"],
                ["가구수", f"{summ.get('household_count'):,}가구" if summ.get("household_count") else "데이터 없음"],
                ["평균 가구원수", f"{summ.get('avg_household_size')}명" if summ.get("avg_household_size") else "데이터 없음"],
            ])
            age = pop.get("age_distribution") or []
            if age:
                doc.add_paragraph("연령 분포")
                _table(["연령대", "인구수"], [[a.get("label", "-"), f"{a.get('count'):,}명" if a.get("count") else "-"] for a in age])
            ht = pop.get("household_types") or []
            if ht:
                doc.add_paragraph("가구원수 분포")
                _table(["가구원수", "비율(%)"], [[h.get("label", "-"), f"{h.get('ratio')}%(추정)" if h.get("ratio") is not None else "-"] for h in ht])
            mig = pop.get("migration") or {}
            doc.add_paragraph("인구 이동")
            _table(["순유입(전입-전출)", "전입", "전출"], [[
                f"{mig.get('net_migration'):,}명" if mig.get("net_migration") is not None else "데이터 없음",
                f"{mig.get('total_inflow'):,}명" if mig.get("total_inflow") is not None else "-",
                f"{mig.get('total_outflow'):,}명" if mig.get("total_outflow") is not None else "-",
            ]])
            _caption(f"출처: {pop.get('source', '-')} ({pop.get('data_source', '-')})")
        else:
            doc.add_paragraph("인구 데이터 없음 / 연동 예정 (인구 분석 미선택)")
        inc = raw.get("income")
        if inc:
            _h("소득 수준", 2)
            avg = inc.get("avg_income_10k")
            med = inc.get("median_income_10k")
            _table(["구분", "연소득"], [
                ["평균 소득", f"{_eok(avg)}원" if avg else "데이터 없음"],
                ["중위 소득" + ("(추정)" if inc.get("median_estimated") else ""), f"{_eok(med)}원" if med else "데이터 없음"],
            ])
            _caption(f"산출근거: {inc.get('source', '-')} ({inc.get('data_source', '-')})")
        else:
            doc.add_paragraph("소득 데이터 없음 / 연동 예정 (소득 분석 미선택)")
        # 타겟 프로파일 요약(주력 연령·가구·소득분위) — 실데이터 축만 정직 표기.
        _tp_bits: list[str] = []
        _pa = tp.get("primary_age") or {}
        _ph = tp.get("primary_household") or {}
        _it = tp.get("income_tier") or {}
        if _pa.get("band"):
            _tp_bits.append(f"주력 연령대 {_pa['band']}")
        if _ph.get("type"):
            _tp_bits.append(f"주력 가구 {_ph['type']}(추정)")
        if _it.get("tier_label"):
            _tp_bits.append(f"소득 수준 {_it['tier_label']}")
        if _tp_bits:
            doc.add_paragraph("타겟 프로파일: " + " · ".join(_tp_bits))
        _ca = _insight_text(rep, "comparable_analysis")
        if _ca:
            doc.add_paragraph(_ca)

        # ── 3. 가격·실거래 (매매·추이·전월세·적정분양가) ──
        _h("3. 가격·실거래", 1)
        _h("매매 시세 (유형별 · 평당가 기준)", 2)
        trade_rows = re_block.get("trade_table") or []
        if trade_rows:
            rows = []
            for s in trade_rows:
                pp = s.get("per_pyeong_manwon")
                area = s.get("avg_area_m2") or 0
                rows.append([
                    s.get("type", "-"),
                    f"{s.get('count', 0)}건",
                    f"{int(pp):,}만원/평" if pp else "-",
                    f"{_eok(s.get('avg_10k', 0))}원" if s.get("avg_10k") else "-",
                    f"{area:.1f}㎡({round(area / PYEONG_SQM)}평)" if area else "-",
                ])
            _table(["유형", "건수", "평당가(만원/평)", "총액 평균", "평균면적"], rows)
        else:
            doc.add_paragraph("데이터 없음")
        _caption(f"출처: {re_block.get('source', '국토교통부 실거래가')} ({re_block.get('data_source', '-')})")

        _h("전월세 보증금 (유형별)", 2)
        rent_rows = re_block.get("rent_table") or []
        if rent_rows:
            rows = [[
                s.get("type", "-"),
                f"{s.get('count', 0)}건",
                f"{_eok(s.get('avg_10k', 0))}원" if s.get("avg_10k") else "-",
                f"{_eok(s.get('min_10k', 0))}원" if s.get("min_10k") else "-",
                f"{_eok(s.get('max_10k', 0))}원" if s.get("max_10k") else "-",
            ] for s in rent_rows]
            _table(["유형", "건수", "평균", "최저", "최고"], rows)
        else:
            doc.add_paragraph("데이터 없음")

        _h("매매 시세 추이 (월별 평당가)", 2)
        trend = [t for t in (re_block.get("trend_series") or []) if t.get("per_pyeong_manwon")]
        if trend:
            rows = [[
                t.get("ym", "-"),
                f"{int(t['per_pyeong_manwon']):,}만원/평" if t.get("per_pyeong_manwon") else "-",
                f"{t['mom_pct']:+.1f}%" if t.get("mom_pct") is not None else "-",
            ] for t in trend]
            _table(["연월", "평당가(만원/평)", "전월대비"], rows)
        else:
            doc.add_paragraph("데이터 없음")

        # 적정 분양가(★사업타당성 앞으로 — 매출/가격이 사업성 산정의 입력)
        _h("적정 분양가 (거래사례비교)", 2)
        pb = rep.get("pricing_band") or (rep.get("analysis") or {}).get("pricing_band") or {}
        if pb and pb.get("data_source") not in (None, "unavailable") and pb.get("fair_price_10k"):
            doc.add_paragraph(f"적정 분양가: {_eok(pb['fair_price_10k'])}원")
            doc.add_paragraph(f"지불여력 판정: {pb.get('affordability_verdict', '-')}")
            _caption(pb.get("note", ""))
        else:
            doc.add_paragraph(pb.get("note") if pb.get("note") else "적정 분양가 산출 불가(비교 데이터 부족) — 가짜값 미생성")
        _pt = _insight_text(rep, "price_trend_analysis", "price_trend")
        if _pt:
            doc.add_paragraph(_pt)

        # ── 4. 입지 분석 (상권 + 인프라) ──
        _h("4. 입지 분석", 1)
        comm = tp.get("commercial") or {}
        loc = tp.get("location") or {}
        if comm.get("data_source") == "live":
            _cat = ", ".join(f"{c.get('category')}({c.get('count')})"
                             for c in (comm.get("category_distribution") or [])[:4])
            doc.add_paragraph(
                f"상권 활성도: {comm.get('grade')}등급(점수 {comm.get('vitality_score')}) · "
                f"점포 {comm.get('total_stores'):,}개" + (f" · 주요업종 {_cat}" if _cat else ""))
        else:
            doc.add_paragraph("상권(SEMAS) 데이터 미확보 — 키 미설정/좌표 부재(정직 표기)")
        if loc.get("nearest_subway"):
            _dist = f" ({loc['subway_distance_m']}m)" if loc.get("subway_distance_m") else ""
            doc.add_paragraph(f"최인접 지하철: {loc['nearest_subway']}{_dist}"
                              + (f" · 학교 {loc['school_count']}개소" if loc.get("school_count") else ""))

        # ── 5. 사업 타당성 (Feasibility · 개략 추정) ──
        _h("5. 사업 타당성 (Feasibility · 개략 추정)", 1)
        fe = rep.get("feasibility_analysis") or {}
        _fin = fe.get("financials") or {}
        _mass = fe.get("massing") or {}
        if _fin and not fe.get("error"):
            _roi = _fin.get("roi_percent")

            def _won(v):
                return f"{_eok(v)}원" if v not in (None, 0) else "-"

            _la = _mass.get("land_area_sqm")
            _gfa = _mass.get("gfa_sqm")
            _table(["구분", "값"], [
                ["ROI(투자수익률)", f"{_roi:.1f}%" if _roi is not None else "-"],
                ["총사업비", _won(_fin.get("total_cost_10k"))],
                ["예상 분양수익", _won(_fin.get("total_revenue_10k"))],
                ["세전 순수익", _won(_fin.get("net_profit_10k"))],
                ["NPV(순현재가치)", _won(_fin.get("npv_10k"))],
                ["대지면적", f"{_la:,.0f}㎡({round(_la / PYEONG_SQM):,}평)" if _la else "-"],
                ["건축가능 연면적", f"{_gfa:,.0f}㎡({round(_mass.get('gfa_pyeong') or 0):,}평)" if _gfa else "-"],
            ])
            _caption((fe.get("assumptions") or {}).get("note", ""))
            um = rep.get("unit_mix_recommendation") or {}
            mix = um.get("recommended_mix") or {}
            if mix and um.get("data_source") not in (None, "unavailable"):
                doc.add_paragraph("권장 평형 배분(수요기반 MD)")
                _table(["평형대", "권장 배분(%)"], [[k, f"{v}%"] for k, v in mix.items()])
                _caption(um.get("rationale", ""))
            elif um.get("data_source") == "unavailable":
                doc.add_paragraph("권장 평형 배분: 데이터 없음(인구/가구 분석 미선택 — 가구원수 분포 필요)")
            _ii = _insight_text(rep, "investment_insight")
            if _ii:
                doc.add_paragraph(_ii)
        else:
            doc.add_paragraph("개략 수지 산출 불가(면적/용도지역 근거 부족) — 무목업 정직 표기")

        # ── 6. 리스크 요인 ──
        _h("6. 리스크 요인", 1)
        _rf = _insight_text(rep, "risk_factors")
        if _rf:
            doc.add_paragraph(_rf)
        _risks = [r for r in (nar.get("risks") or []) if str(r).strip()]
        for r in _risks:
            doc.add_paragraph(r, style="List Bullet")
        if not _rf and not _risks:
            doc.add_paragraph(_ai_missing)

        # ── 7. 결론 및 권고 ──
        _h("7. 결론 및 권고", 1)
        doc.add_paragraph(f"투자의견: {es['opinion']}")
        _tr = _insight_text(rep, "timing_recommendation")
        if _tr:
            doc.add_paragraph(_tr)
        _ii2 = _insight_text(rep, "investment_insight")
        if _ii2:
            doc.add_paragraph(_ii2)
        _opps = [o for o in (nar.get("opportunities") or []) if str(o).strip()]
        if _opps:
            doc.add_paragraph("기회 요인")
            for o in _opps:
                doc.add_paragraph(o, style="List Bullet")
        _persona = nar.get("target_persona") or (rep.get("analysis") or {}).get("target_persona")
        if _persona and _persona not in ("AI 분석 미포함",):
            doc.add_paragraph(f"추천 분양 타겟(페르소나): {_persona}")

        # ── 8. 부록 · 출처 및 면책 ──
        _h("8. 부록 · 출처 및 면책", 1)
        doc.add_paragraph(f"실거래 출처: {re_block.get('source', '국토교통부 실거래가')}")
        prem = tp.get("premium") or {}
        if prem.get("note"):
            _caption(prem["note"])
        _caption(DISCLAIMER_TEXT)

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()
